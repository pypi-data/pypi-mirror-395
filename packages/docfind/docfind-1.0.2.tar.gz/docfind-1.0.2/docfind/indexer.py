"""
Document indexing module.

Handles extraction and indexing of text from various document formats.
Supports PDF, DOCX, XLSX, PPTX, HTML, text files, and unknown formats via hex extraction.
"""

import logging
import concurrent.futures
from typing import Optional, Callable, Dict, Any, List, Tuple
from pathlib import Path
from datetime import datetime
import traceback

from .db import DatabaseManager
from .utils import (
    calculate_sha256,
    detect_file_type,
    is_text_file,
    should_ignore,
    format_size,
)
from .hex_extractor import extract_text_from_unknown

logger = logging.getLogger(__name__)


class TextExtractor:
    """Extract text from various document formats."""

    def __init__(self, trust_external_tools: bool = False):
        """
        Initialize text extractor.

        Args:
            trust_external_tools: Whether to use external conversion tools
        """
        self.trust_external_tools = trust_external_tools

    def extract(self, file_path: Path, file_type: str) -> Tuple[str, str]:
        """
        Extract text from a file.

        Args:
            file_path: Path to file
            file_type: File type/extension

        Returns:
            Tuple of (extracted_text, extractor_name)
        """
        # Map file types to extraction methods
        extractors = {
            "pdf": self._extract_pdf,
            "docx": self._extract_docx,
            "doc": self._extract_doc,
            "xlsx": self._extract_xlsx,
            "xls": self._extract_xls,
            "pptx": self._extract_pptx,
            "ppt": self._extract_ppt,
            "html": self._extract_html,
            "htm": self._extract_html,
            "xml": self._extract_xml,
            "txt": self._extract_text,
            "md": self._extract_text,
            "rst": self._extract_text,
            "csv": self._extract_text,
            "json": self._extract_text,
            "log": self._extract_text,
            "py": self._extract_text,
            "js": self._extract_text,
            "java": self._extract_text,
            "c": self._extract_text,
            "cpp": self._extract_text,
            "h": self._extract_text,
            "cs": self._extract_text,
            "go": self._extract_text,
            "rs": self._extract_text,
            "rb": self._extract_text,
            "php": self._extract_text,
            "sh": self._extract_text,
            "bat": self._extract_text,
            "ps1": self._extract_text,
        }

        extractor = extractors.get(file_type, self._extract_unknown)
        return extractor(file_path)

    def _extract_pdf(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from PDF."""
        try:
            from pdfminer.high_level import extract_text

            text = extract_text(str(file_path))
            return text, "pdfminer"
        except Exception as e:
            logger.error(f"PDF extraction failed for {file_path}: {e}")
            # Fallback to hex extraction
            return self._extract_unknown(file_path)

    def _extract_docx(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from DOCX."""
        try:
            from docx import Document

            doc = Document(str(file_path))
            text = "\n".join([para.text for para in doc.paragraphs])
            return text, "python-docx"
        except Exception as e:
            logger.error(f"DOCX extraction failed for {file_path}: {e}")
            return self._extract_unknown(file_path)

    def _extract_doc(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from DOC (old Word format)."""
        # Old .doc format requires external tools
        logger.warning(
            f"Old .doc format not supported, using hex extraction: {file_path}"
        )
        return self._extract_unknown(file_path)

    def _extract_xlsx(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from XLSX."""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(str(file_path), read_only=True, data_only=True)

            text_parts = []
            for sheet in wb.worksheets:
                text_parts.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(
                        [str(cell) if cell is not None else "" for cell in row]
                    )
                    if row_text.strip():
                        text_parts.append(row_text)

            return "\n".join(text_parts), "openpyxl"

        except Exception as e:
            logger.error(f"XLSX extraction failed for {file_path}: {e}")
            return self._extract_unknown(file_path)

    def _extract_xls(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from XLS (old Excel format)."""
        logger.warning(
            f"Old .xls format not supported, using hex extraction: {file_path}"
        )
        return self._extract_unknown(file_path)

    def _extract_pptx(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from PPTX."""
        try:
            from pptx import Presentation

            prs = Presentation(str(file_path))

            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"[Slide {slide_num}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            text_parts.append(shape.text)

            return "\n".join(text_parts), "python-pptx"

        except Exception as e:
            logger.error(f"PPTX extraction failed for {file_path}: {e}")
            return self._extract_unknown(file_path)

    def _extract_ppt(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from PPT (old PowerPoint format)."""
        logger.warning(
            f"Old .ppt format not supported, using hex extraction: {file_path}"
        )
        return self._extract_unknown(file_path)

    def _extract_html(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from HTML."""
        try:
            from bs4 import BeautifulSoup

            with open(file_path, "rb") as f:
                soup = BeautifulSoup(f, "lxml")

            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()

            text = soup.get_text(separator="\n", strip=True)
            return text, "beautifulsoup4"

        except Exception as e:
            logger.error(f"HTML extraction failed for {file_path}: {e}")
            return self._extract_text(file_path)

    def _extract_xml(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from XML."""
        try:
            from bs4 import BeautifulSoup

            with open(file_path, "rb") as f:
                soup = BeautifulSoup(f, "lxml-xml")

            text = soup.get_text(separator="\n", strip=True)
            return text, "beautifulsoup4"

        except Exception as e:
            logger.error(f"XML extraction failed for {file_path}: {e}")
            return self._extract_text(file_path)

    def _extract_text(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from plain text file."""
        try:
            # Try multiple encodings
            for encoding in ["utf-8", "utf-8-sig", "latin-1", "cp1252"]:
                try:
                    with open(file_path, "r", encoding=encoding) as f:
                        text = f.read()
                    return text, f"text ({encoding})"
                except UnicodeDecodeError:
                    continue

            # If all encodings fail, use binary mode with error handling
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
            return text, "text (utf-8 with errors)"

        except Exception as e:
            logger.error(f"Text extraction failed for {file_path}: {e}")
            return "", "error"

    def _extract_unknown(self, file_path: Path) -> Tuple[str, str]:
        """Extract text from unknown format using hex extraction."""
        try:
            text = extract_text_from_unknown(file_path)
            return text, "hex-extractor"
        except Exception as e:
            logger.error(f"Hex extraction failed for {file_path}: {e}")
            return "", "error"


class IndexProgress:
    """Progress tracking for indexing operations."""

    def __init__(self):
        self.total_files = 0
        self.processed_files = 0
        self.successful = 0
        self.failed = 0
        self.skipped = 0
        self.bytes_processed = 0
        self.current_file = ""
        self.start_time = datetime.now()
        self.errors: List[Dict[str, str]] = []

    def update(
        self,
        file_path: str,
        status: str,
        file_size: int = 0,
        error: Optional[str] = None,
    ):
        """Update progress."""
        self.processed_files += 1
        self.current_file = file_path
        self.bytes_processed += file_size

        if status == "success":
            self.successful += 1
        elif status == "error":
            self.failed += 1
            if error:
                self.errors.append({"file": file_path, "error": error})
        elif status == "skipped":
            self.skipped += 1

    def get_stats(self) -> Dict[str, Any]:
        """Get progress statistics."""
        elapsed = (datetime.now() - self.start_time).total_seconds()

        return {
            "total_files": self.total_files,
            "processed_files": self.processed_files,
            "successful": self.successful,
            "failed": self.failed,
            "skipped": self.skipped,
            "bytes_processed": self.bytes_processed,
            "current_file": self.current_file,
            "elapsed_seconds": elapsed,
            "files_per_second": self.processed_files / elapsed if elapsed > 0 else 0,
            "errors": self.errors[-10:],  # Last 10 errors
        }


class DocumentIndexer:
    """Index documents into the database."""

    def __init__(
        self,
        db: DatabaseManager,
        max_file_size: int = 100 * 1024 * 1024,
        trust_external_tools: bool = False,
        ignore_patterns: Optional[List[str]] = None,
    ):
        """
        Initialize document indexer.

        Args:
            db: Database manager
            max_file_size: Maximum file size to index (bytes)
            trust_external_tools: Whether to use external tools
            ignore_patterns: List of glob patterns to ignore
        """
        self.db = db
        self.max_file_size = max_file_size
        self.trust_external_tools = trust_external_tools
        self.ignore_patterns = ignore_patterns or []
        self.extractor = TextExtractor(trust_external_tools)

    def index_directory(
        self,
        root_path: Path,
        reindex: bool = False,
        threads: int = 4,
        progress_callback: Optional[Callable[[IndexProgress], None]] = None,
    ) -> IndexProgress:
        """
        Index all files in a directory.

        Args:
            root_path: Root directory to index
            reindex: Whether to reindex existing files
            threads: Number of worker threads
            progress_callback: Optional callback for progress updates

        Returns:
            IndexProgress object with statistics
        """
        root_path = root_path.resolve()
        logger.info(f"Indexing directory: {root_path}")

        if not root_path.exists():
            raise ValueError(f"Path does not exist: {root_path}")

        if not root_path.is_dir():
            raise ValueError(f"Path is not a directory: {root_path}")

        # Collect files to index
        files = self._collect_files(root_path)

        progress = IndexProgress()
        progress.total_files = len(files)

        logger.info(f"Found {len(files)} files to index")

        if reindex:
            logger.info("Removing existing entries for this root...")
            self.db.delete_by_root(str(root_path))

        # Process files
        if threads == 1:
            # Single-threaded
            for file_path in files:
                self._index_file(file_path, root_path, progress, progress_callback)
        else:
            # Multi-threaded
            with concurrent.futures.ThreadPoolExecutor(max_workers=threads) as executor:
                futures = [
                    executor.submit(
                        self._index_file,
                        file_path,
                        root_path,
                        progress,
                        progress_callback,
                    )
                    for file_path in files
                ]

                concurrent.futures.wait(futures)

        logger.info(
            f"Indexing complete: {progress.successful} successful, {progress.failed} failed, {progress.skipped} skipped"
        )

        return progress

    def _collect_files(self, root_path: Path) -> List[Path]:
        """Collect all files to index."""
        files = []

        for path in root_path.rglob("*"):
            if not path.is_file():
                continue

            # Check ignore patterns
            if should_ignore(path, self.ignore_patterns):
                logger.debug(f"Ignoring {path} (matches ignore pattern)")
                continue

            # Check file size
            try:
                size = path.stat().st_size
                if size > self.max_file_size:
                    logger.debug(f"Ignoring {path} (too large: {format_size(size)})")
                    continue
            except Exception as e:
                logger.warning(f"Failed to stat {path}: {e}")
                continue

            files.append(path)

        return files

    def _index_file(
        self,
        file_path: Path,
        root_path: Path,
        progress: IndexProgress,
        progress_callback: Optional[Callable[[IndexProgress], None]],
    ):
        """Index a single file."""
        try:
            # Get file metadata
            stat = file_path.stat()
            file_size = stat.st_size
            mtime = stat.st_mtime

            # Detect file type
            file_type = detect_file_type(file_path)

            # Check if file needs reindexing
            existing = self.db.get_document_by_path(str(file_path))
            if existing and existing.get("mtime") == mtime:
                logger.debug(f"Skipping {file_path} (not modified)")
                progress.update(str(file_path), "skipped", file_size)
                if progress_callback:
                    progress_callback(progress)
                return

            # Extract text
            try:
                extracted_text, extractor_name = self.extractor.extract(
                    file_path, file_type
                )
            except Exception as e:
                logger.error(f"Extraction failed for {file_path}: {e}")
                extracted_text = ""
                extractor_name = "error"

            # Calculate hash
            try:
                sha256 = calculate_sha256(file_path)
            except Exception as e:
                logger.warning(f"Failed to calculate hash for {file_path}: {e}")
                sha256 = None

            # Insert into database
            status = "success" if extracted_text else "error"
            error_message = None if status == "success" else "No text extracted"

            self.db.insert_document(
                path=str(file_path),
                root_path=str(root_path),
                source_type=file_type,
                extracted_text=extracted_text,
                extractor=extractor_name,
                sha256=sha256,
                file_size=file_size,
                mtime=mtime,
                status=status,
                error_message=error_message,
            )

            progress.update(str(file_path), status, file_size)
            logger.debug(f"Indexed {file_path} ({extractor_name})")

        except Exception as e:
            error_msg = f"{type(e).__name__}: {str(e)}"
            logger.error(f"Failed to index {file_path}: {error_msg}")
            logger.debug(traceback.format_exc())

            # Try to insert error record
            try:
                self.db.insert_document(
                    path=str(file_path),
                    root_path=str(root_path),
                    source_type=detect_file_type(file_path),
                    extracted_text="",
                    status="error",
                    error_message=error_msg,
                )
            except Exception as db_error:
                logger.error(f"Failed to record error for {file_path}: {db_error}")

            progress.update(str(file_path), "error", 0, error_msg)

        finally:
            if progress_callback:
                progress_callback(progress)

    def index_single_file(
        self, file_path: Path, root_path: Optional[Path] = None
    ) -> bool:
        """
        Index a single file.

        Args:
            file_path: Path to file
            root_path: Optional root path (defaults to file's parent)

        Returns:
            True if successful
        """
        if root_path is None:
            root_path = file_path.parent

        progress = IndexProgress()
        progress.total_files = 1

        self._index_file(file_path, root_path, progress, None)

        return progress.successful > 0
