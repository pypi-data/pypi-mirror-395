import base64
import json
import logging
import re
from dataclasses import dataclass
from enum import Enum
from pathlib import Path
from typing import Dict, List, Any, Union, Optional, Tuple


class TableStyle(Enum):
    """Table output style options for complex tables with merged cells"""
    MINIMAL_HTML = "minimal_html"      # Clean HTML with only rowspan/colspan, no inline styles
    MARKDOWN_GRID = "markdown_grid"    # Extended markdown with merge annotations
    STYLED_HTML = "styled_html"        # Full HTML with inline styles (legacy)
    
    @classmethod
    def default(cls):
        return cls.MINIMAL_HTML

# Office document libraries
try:
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    from docx.shape import InlineShape
    from docx.enum.shape import WD_INLINE_SHAPE
    from docx.oxml.ns import qn
    import docx.oxml.text.paragraph
    import docx.oxml.table
except ImportError:
    raise ImportError("python-docx is required. Install with: pip install python-docx")

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    raise ImportError("python-pptx is required. Install with: pip install python-pptx")

try:
    import openpyxl
    from openpyxl.drawing.image import Image as XLImage
except ImportError:
    raise ImportError("openpyxl is required. Install with: pip install openpyxl")

try:
    import pandas as pd
except ImportError:
    raise ImportError("pandas is required. Install with: pip install pandas")

# Import VisionAgent for OCR functionality (optional)
try:
    from doc2mark.ocr.openai import VisionAgent

    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    logging.warning("OCR functionality not available. Install VisionAgent to enable OCR.")

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


@dataclass
class SimpleContent:
    """Simple content item with type and data"""
    type: str  # 'text:title', 'text:section', 'text:normal', 'text:list', 'text:caption', 'table', or 'image'
    content: str  # markdown text, markdown table, or base64 data
    page: int  # For Office docs, this might be slide number or sheet number
    position_y: float  # For sorting (approximate position)


class BaseOfficeLoader:
    """Base class for Office document loaders"""

    def __init__(self, file_path: Union[str, Path], ocr=None, table_style: Union[str, TableStyle] = None):
        self.file_path = Path(file_path)
        if not self.file_path.exists():
            raise FileNotFoundError(f"File not found: {self.file_path}")
        self.doc = None
        self.ocr = ocr  # Store the OCR instance
        
        # Set table output style
        if table_style is None:
            self.table_style = TableStyle.default()
        elif isinstance(table_style, str):
            self.table_style = TableStyle(table_style)
        else:
            self.table_style = table_style

        # Log OCR configuration if available
        if self.ocr:
            logger.info(f"📷 OCR configured for {self.__class__.__name__}: {type(self.ocr).__name__}")
            if hasattr(self.ocr, 'config') and self.ocr.config:
                if hasattr(self.ocr.config, 'language') and self.ocr.config.language:
                    logger.info(f"🌍 OCR Language setting: {self.ocr.config.language}")
        else:
            logger.warning(f"⚠️  No OCR instance provided to {self.__class__.__name__}")

    def convert_to_json(self,
                        extract_images: bool = True,
                        ocr_images: bool = False,
                        show_progress: bool = True) -> Dict[str, Any]:
        """Convert document to JSON format"""
        raise NotImplementedError("Subclasses must implement this method")

    def _convert_table_to_markdown(self, table_data: Union[List[List[str]], Any],
                                   extract_images: bool = True, ocr_images: bool = False,
                                   ocr_results_map: Dict[str, str] = {}) -> str:
        """Convert table data to markdown format with complex structure support
        
        This enhanced version handles:
        - Merged cells detection
        - Complex table structures
        - Automatic HTML conversion for complex tables
        - Line breaks preservation
        - Images within cells (with OCR support)
        
        Library-specific merged cell handling:
        - DOCX: Merged cells share the same cell object (_tc), detected by comparing object IDs
        - PPTX: Merged cells have empty strings '' in non-origin cells, with gridSpan/vMerge properties
        - XLSX: Merged cells have None in non-origin cells, with explicit merged_cells.ranges
        """
        # Handle different table input types
        if hasattr(table_data, 'rows'):  # DOCX Table object
            rows_data = []
            merged_cells_info = []

            # Get table dimensions
            num_rows = len(table_data.rows)
            num_cols = len(table_data.columns)

            logger.debug(f"DOCX Table dimensions: {num_rows}x{num_cols}")

            # Initialize a 2D array to track cell content and merged status
            table_array = [[None for _ in range(num_cols)] for _ in range(num_rows)]
            cell_map = {}  # Maps cell objects to their content and span info
            processed_positions = set()

            # Process each cell by coordinates
            for row_idx in range(num_rows):
                logger.debug(f"\nProcessing row {row_idx}:")
                for col_idx in range(num_cols):
                    if (row_idx, col_idx) in processed_positions:
                        logger.debug(f"  Cell ({row_idx}, {col_idx}): Already processed (part of merge)")
                        continue

                    try:
                        # Access cell by coordinates
                        cell = table_data.cell(row_idx, col_idx)
                        cell_id = id(cell._tc) if hasattr(cell, '_tc') else id(cell)

                        # Get the actual span of this cell by checking adjacent cells
                        min_row, min_col = row_idx, col_idx
                        max_row, max_col = row_idx, col_idx

                        # Find the full extent of this merged cell by checking all directions
                        # We need to find the rectangle that encompasses all positions with the same cell ID

                        # First, find all positions that share this cell ID
                        cell_positions = [(row_idx, col_idx)]

                        # Check all remaining positions in the table
                        for r in range(num_rows):
                            for c in range(num_cols):
                                if (r, c) != (row_idx, col_idx) and (r, c) not in processed_positions:
                                    try:
                                        test_cell = table_data.cell(r, c)
                                        test_id = id(test_cell._tc) if hasattr(test_cell, '_tc') else id(test_cell)
                                        if test_id == cell_id:
                                            cell_positions.append((r, c))
                                            min_row = min(min_row, r)
                                            max_row = max(max_row, r)
                                            min_col = min(min_col, c)
                                            max_col = max(max_col, c)
                                    except:
                                        pass

                        # Extract cell content including images with OCR
                        cell_text = self._extract_cell_content_with_images(
                            cell, extract_images, ocr_images, ocr_results_map
                        )

                        logger.debug(
                            f"  Cell ({row_idx}, {col_idx}): '{cell_text[:30]}...' spans to ({max_row}, {max_col})")

                        # Place text in the array
                        table_array[min_row][min_col] = cell_text

                        # Mark all positions covered by this cell
                        for r in range(min_row, max_row + 1):
                            for c in range(min_col, max_col + 1):
                                processed_positions.add((r, c))
                                if (r, c) != (min_row, min_col):
                                    table_array[r][c] = ""  # Empty string for merged areas
                                    logger.debug(f"    Marking ({r}, {c}) as empty (part of merge)")

                        # Record merge info if this is a merged cell
                        rowspan = max_row - min_row + 1
                        colspan = max_col - min_col + 1
                        if rowspan > 1 or colspan > 1:
                            merged_cells_info.append({
                                'row': min_row,
                                'col': min_col,
                                'rowspan': rowspan,
                                'colspan': colspan
                            })
                            logger.debug(
                                f"Merged cell at ({min_row}, {min_col}): {rowspan}x{colspan} - '{cell_text[:30]}...'")

                    except Exception as e:
                        logger.warning(f"Error processing cell at ({row_idx}, {col_idx}): {e}")
                        table_array[row_idx][col_idx] = ""

            # Log the final table array for debugging
            logger.debug("\nFinal table array:")
            for i, row in enumerate(table_array):
                logger.debug(f"  Row {i}: {[cell[:20] + '...' if cell and len(cell) > 20 else cell for cell in row]}")

            # Convert array to list format
            for row in table_array:
                rows_data.append([cell if cell is not None else "" for cell in row])

            table_data = rows_data
            is_complex = len(merged_cells_info) > 0

            logger.info(f"DOCX Table: Detected {len(merged_cells_info)} merged cells, complex={is_complex}")

            # If complex, create table info for HTML conversion
            if is_complex:
                table_info = {
                    'is_complex': True,
                    'merged_cells': merged_cells_info,
                    'row_count': num_rows,
                    'col_count': num_cols,
                    'cell_spans': {}
                }

                # Add span info
                for merge_info in merged_cells_info:
                    key = (merge_info['row'], merge_info['col'])
                    table_info['cell_spans'][key] = (merge_info['rowspan'], merge_info['colspan'])

                return self._convert_table_to_html(table_data, table_info)
            else:
                # For DOCX tables without merges, use simple markdown directly
                # Don't call _analyze_table_structure which might misinterpret empty cells
                table_info = {
                    'is_complex': False,
                    'merged_cells': [],
                    'row_count': num_rows,
                    'col_count': num_cols,
                    'cell_spans': {}
                }
                return self._convert_table_to_simple_markdown(table_data, table_info)
        else:
            # For non-DOCX tables (e.g., from PPTX or XLSX), still analyze structure
            is_complex = False

        if not table_data or not any(table_data):
            return ""

        # Only analyze table structure for non-DOCX tables
        table_info = self._analyze_table_structure(table_data)

        # If table is complex, use HTML format
        if is_complex or table_info['is_complex']:
            return self._convert_table_to_html(table_data, table_info)
        else:
            return self._convert_table_to_simple_markdown(table_data, table_info)

    def _analyze_table_structure(self, table_data: List[List]) -> Dict[str, Any]:
        """Analyze table structure to detect merged cells and complexity
        
        This method is used for non-DOCX tables (PPTX and XLSX) where we need to detect merges
        based on empty cells. DOCX tables handle merges differently using cell object IDs.
        """
        if not table_data:
            return {'is_complex': False, 'merged_cells': [], 'row_count': 0, 'col_count': 0, 'cell_spans': {}}

        row_count = len(table_data)
        col_count = max(len(row) for row in table_data) if table_data else 0

        # Initialize analysis structures
        cell_spans = {}
        merged_cells = []
        is_complex = False

        # Create a normalized table (all rows same length)
        normalized = []
        for row in table_data:
            normalized_row = list(row) + [""] * (col_count - len(row))
            normalized.append(normalized_row)

        # Detect merged cells by looking for patterns
        for row_idx in range(row_count):
            for col_idx in range(col_count):
                cell = normalized[row_idx][col_idx]

                if cell == "":
                    # Check if this is part of a merged cell
                    span_info = self._detect_cell_span(normalized, row_idx, col_idx)
                    if span_info:
                        merged_cells.append(span_info)
                        is_complex = True
                else:
                    # Check if this cell spans multiple rows/cols
                    rowspan, colspan = self._calculate_cell_span(normalized, row_idx, col_idx, str(cell))
                    if rowspan > 1 or colspan > 1:
                        cell_spans[(row_idx, col_idx)] = (rowspan, colspan)
                        merged_cells.append({
                            'row': row_idx,
                            'col': col_idx,
                            'rowspan': rowspan,
                            'colspan': colspan,
                            'content': str(cell)
                        })
                        is_complex = True

        return {
            'is_complex': is_complex,
            'merged_cells': merged_cells,
            'row_count': row_count,
            'col_count': col_count,
            'cell_spans': cell_spans
        }

    def _detect_cell_span(self, table: List[List], row: int, col: int) -> Optional[Dict]:
        """Detect if an empty cell is part of a span from another cell
        
        Empty cells are represented as:
        - None in XLSX (openpyxl)
        - '' (empty string) in PPTX (python-pptx)
        """
        cell_value = table[row][col]
        
        # Check if this cell is truly empty (None or empty string)
        if not (cell_value is None or (isinstance(cell_value, str) and cell_value == '')):
            return None
            
        # Check if empty cell is part of a row span from above
        if row > 0:
            above_cell = table[row - 1][col]
            # Check if cell above has content
            if above_cell is not None and str(above_cell).strip() != '':
                # Check if cells below are also empty (indicating rowspan)
                span_rows = 1
                for check_row in range(row, len(table)):
                    check_cell = table[check_row][col]
                    if check_cell is None or (isinstance(check_cell, str) and check_cell == ''):
                        span_rows += 1
                    else:
                        break

                if span_rows > 1:
                    return {
                        'type': 'rowspan_continuation',
                        'source_row': row - 1,
                        'source_col': col,
                        'span_rows': span_rows
                    }

        # Check if empty cell is part of a col span from left
        if col > 0:
            left_cell = table[row][col - 1]
            # Check if cell to the left has content
            if left_cell is not None and str(left_cell).strip() != '':
                # Check if cells to right are also empty (indicating colspan)
                span_cols = 1
                for check_col in range(col, len(table[row])):
                    check_cell = table[row][check_col]
                    if check_cell is None or (isinstance(check_cell, str) and check_cell == ''):
                        span_cols += 1
                    else:
                        break

                if span_cols > 1:
                    return {
                        'type': 'colspan_continuation',
                        'source_row': row,
                        'source_col': col - 1,
                        'span_cols': span_cols
                    }

        return None

    def _calculate_cell_span(self, table: List[List], row: int, col: int, cell_content: str) -> Tuple[int, int]:
        """Calculate how many rows and columns a cell spans
        
        Only counts consecutive None or empty string cells as part of the span
        """
        rowspan = 1
        colspan = 1

        # Check colspan: count consecutive empty cells to the right
        for check_col in range(col + 1, len(table[row])):
            check_cell = table[row][check_col]
            if check_cell is None or (isinstance(check_cell, str) and check_cell == ''):
                colspan += 1
            else:
                break

        # Check rowspan: count consecutive empty cells below
        for check_row in range(row + 1, len(table)):
            if col < len(table[check_row]):
                check_cell = table[check_row][col]
                if check_cell is None or (isinstance(check_cell, str) and check_cell == ''):
                    rowspan += 1
                else:
                    break
            else:
                break

        return rowspan, colspan

    def _convert_table_to_simple_markdown(self, table_data: List[List], table_info: Dict) -> str:
        """Convert simple table to markdown format with span annotations"""
        if not table_data:
            return ""

        markdown_lines = []
        processed_cells = set()  # Track cells that are part of spans

        # Add table complexity note if needed
        if table_info['merged_cells']:
            markdown_lines.append("<!-- Table contains merged cells (marked with *) -->")

        # Build markdown table
        for row_idx, row in enumerate(table_data):
            row_cells = []

            for col_idx in range(table_info['col_count']):
                # Skip if this cell is part of a span
                if (row_idx, col_idx) in processed_cells:
                    continue

                # Get cell content
                if col_idx < len(row):
                    cell_text = str(row[col_idx]).strip()
                else:
                    cell_text = ""

                # Replace newlines with <br> tags
                cell_text = "<br>".join(cell_text.split('\n'))
                # Escape pipe characters
                cell_text = cell_text.replace("|", "\\|")

                # Check if this cell has spans
                if (row_idx, col_idx) in table_info['cell_spans']:
                    rowspan, colspan = table_info['cell_spans'][(row_idx, col_idx)]

                    # Mark spanned cells as processed
                    for r in range(row_idx, min(row_idx + rowspan, table_info['row_count'])):
                        for c in range(col_idx, min(col_idx + colspan, table_info['col_count'])):
                            if r != row_idx or c != col_idx:
                                processed_cells.add((r, c))

                    # Add span indicator
                    if rowspan > 1 or colspan > 1:
                        span_note = f"*[{rowspan}x{colspan}]*"
                        cell_text = f"{cell_text} {span_note}" if cell_text else span_note

                # For cells that span multiple columns, repeat the content
                if (row_idx, col_idx) in table_info['cell_spans']:
                    _, colspan = table_info['cell_spans'][(row_idx, col_idx)]
                    for _ in range(colspan):
                        row_cells.append(cell_text)
                else:
                    row_cells.append(cell_text)

            # Ensure row has correct number of columns
            while len(row_cells) < table_info['col_count']:
                row_cells.append("")

            # Create table row
            row_text = "| " + " | ".join(row_cells[:table_info['col_count']]) + " |"
            markdown_lines.append(row_text)

            # Add separator after first row
            if row_idx == 0:
                separator = "|" + "|".join([" --- " for _ in range(table_info['col_count'])]) + "|"
                markdown_lines.append(separator)

        return "\n".join(markdown_lines) + "\n\n"

    def _convert_table_to_html(self, table_data: List[List], table_info: Dict) -> str:
        """Convert complex table to HTML format based on table_style setting"""
        if not table_data:
            return ""
        
        # Route to appropriate converter based on style
        if self.table_style == TableStyle.STYLED_HTML:
            return self._convert_table_to_styled_html(table_data, table_info)
        elif self.table_style == TableStyle.MARKDOWN_GRID:
            return self._convert_table_to_markdown_grid(table_data, table_info)
        else:  # MINIMAL_HTML (default)
            return self._convert_table_to_minimal_html(table_data, table_info)
    
    def _convert_table_to_minimal_html(self, table_data: List[List], table_info: Dict) -> str:
        """Convert complex table to clean, minimal HTML without inline styles"""
        if not table_data:
            return ""

        html_lines = ["<table>"]
        processed_cells = set()

        for row_idx, row in enumerate(table_data):
            html_lines.append("<tr>")

            col_idx = 0
            while col_idx < table_info['col_count']:
                if (row_idx, col_idx) in processed_cells and (row_idx, col_idx) not in table_info['cell_spans']:
                    col_idx += 1
                    continue

                # Get and escape cell content
                cell_text = str(row[col_idx]).strip() if col_idx < len(row) and row[col_idx] is not None else ""
                cell_text = cell_text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                cell_text = cell_text.replace('\n', '<br>')

                # Build attributes (only rowspan/colspan if needed)
                attrs = []
                colspan = 1
                rowspan = 1

                if (row_idx, col_idx) in table_info['cell_spans']:
                    rowspan, colspan = table_info['cell_spans'][(row_idx, col_idx)]
                    if rowspan > 1:
                        attrs.append(f'rowspan="{rowspan}"')
                    if colspan > 1:
                        attrs.append(f'colspan="{colspan}"')
                    
                    for r in range(row_idx, min(row_idx + rowspan, table_info['row_count'])):
                        for c in range(col_idx, min(col_idx + colspan, table_info['col_count'])):
                            processed_cells.add((r, c))

                cell_tag = "th" if row_idx == 0 else "td"
                attrs_str = " " + " ".join(attrs) if attrs else ""
                html_lines.append(f"<{cell_tag}{attrs_str}>{cell_text}</{cell_tag}>")

                col_idx += colspan

            html_lines.append("</tr>")

        html_lines.append("</table>")
        return "\n".join(html_lines) + "\n\n"
    
    def _convert_table_to_markdown_grid(self, table_data: List[List], table_info: Dict) -> str:
        """Convert complex table to markdown with merge annotations as comments"""
        if not table_data:
            return ""
        
        lines = []
        processed_cells = set()
        
        # Add merge info as a compact header comment
        if table_info.get('cell_spans'):
            merge_notes = []
            for (r, c), (rowspan, colspan) in table_info['cell_spans'].items():
                if rowspan > 1 or colspan > 1:
                    merge_notes.append(f"R{r+1}C{c+1}:{rowspan}x{colspan}")
            if merge_notes:
                lines.append(f"<!-- Merged: {', '.join(merge_notes)} -->")
        
        # Calculate column widths (no truncation to preserve data integrity)
        col_count = table_info['col_count']
        col_widths = [3] * col_count  # minimum width
        
        for row in table_data:
            for i, cell in enumerate(row[:col_count]):
                cell_text = str(cell).strip() if cell else ""
                col_widths[i] = max(col_widths[i], len(cell_text))
        
        # Build markdown table
        for row_idx, row in enumerate(table_data):
            row_cells = []
            col_idx = 0
            
            while col_idx < col_count:
                if (row_idx, col_idx) in processed_cells and (row_idx, col_idx) not in table_info.get('cell_spans', {}):
                    # Cell consumed by span - use empty placeholder
                    row_cells.append("↓" if any(
                        r < row_idx and c == col_idx and (r, c) in table_info.get('cell_spans', {})
                        for r in range(row_idx) for c in [col_idx]
                    ) else "→")
                    col_idx += 1
                    continue
                
                cell_text = str(row[col_idx]).strip() if col_idx < len(row) and row[col_idx] is not None else ""
                
                # Mark merged cells with indicator
                if (row_idx, col_idx) in table_info.get('cell_spans', {}):
                    rowspan, colspan = table_info['cell_spans'][(row_idx, col_idx)]
                    if rowspan > 1 or colspan > 1:
                        cell_text = f"{cell_text} ⊕" if cell_text else "⊕"
                    
                    for r in range(row_idx, min(row_idx + rowspan, table_info['row_count'])):
                        for c in range(col_idx, min(col_idx + colspan, col_count)):
                            processed_cells.add((r, c))
                
                row_cells.append(cell_text.ljust(col_widths[col_idx]))
                col_idx += 1
            
            lines.append("| " + " | ".join(row_cells) + " |")
            
            # Add separator after header
            if row_idx == 0:
                sep_cells = ["-" * w for w in col_widths]
                lines.append("| " + " | ".join(sep_cells) + " |")
        
        return "\n".join(lines) + "\n\n"
    
    def _convert_table_to_styled_html(self, table_data: List[List], table_info: Dict) -> str:
        """Convert complex table to HTML with full inline styles (legacy format)"""
        if not table_data:
            return ""

        html_lines = ["<!-- Complex table converted to HTML for better structure preservation -->"]
        html_lines.append('<table border="1" style="border-collapse: collapse; width: 100%;">')

        processed_cells = set()

        # Process each row
        for row_idx, row in enumerate(table_data):
            html_lines.append("  <tr>")

            col_idx = 0
            while col_idx < table_info['col_count']:
                # Skip if this cell is part of a vertical span from a previous row
                if (row_idx, col_idx) in processed_cells and (row_idx, col_idx) not in table_info['cell_spans']:
                    # This cell is consumed by a rowspan from above, skip it
                    col_idx += 1
                    continue

                # Get cell content
                if col_idx < len(row) and row[col_idx] is not None:
                    cell_text = str(row[col_idx]).strip()
                else:
                    cell_text = ""

                # Convert newlines to <br> for HTML
                cell_text = cell_text.replace('\n', '<br>')

                # Escape HTML special characters
                cell_text = cell_text.replace('&', '&amp;')
                cell_text = cell_text.replace('<', '&lt;')
                cell_text = cell_text.replace('>', '&gt;')
                cell_text = cell_text.replace('"', '&quot;')

                # But restore <br> tags we just added
                cell_text = cell_text.replace('&lt;br&gt;', '<br>')

                # Determine cell attributes
                cell_attrs = []
                cell_style = []
                colspan = 1
                rowspan = 1

                # Check if this cell has spans
                if (row_idx, col_idx) in table_info['cell_spans']:
                    rowspan, colspan = table_info['cell_spans'][(row_idx, col_idx)]

                    if rowspan > 1:
                        cell_attrs.append(f'rowspan="{rowspan}"')
                    if colspan > 1:
                        cell_attrs.append(f'colspan="{colspan}"')

                    # Mark spanned cells as processed
                    for r in range(row_idx, min(row_idx + rowspan, table_info['row_count'])):
                        for c in range(col_idx, min(col_idx + colspan, table_info['col_count'])):
                            processed_cells.add((r, c))

                # Determine if header cell (first row typically)
                cell_tag = "th" if row_idx == 0 else "td"

                # Add some basic styling
                if cell_tag == "th":
                    cell_style.append("background-color: #f0f0f0")
                    cell_style.append("font-weight: bold")

                cell_style.append("padding: 8px")
                cell_style.append("text-align: left")
                cell_style.append("vertical-align: top")
                cell_style.append("border: 1px solid #ddd")

                # Build style attribute
                if cell_style:
                    cell_attrs.append(f'style="{"; ".join(cell_style)}"')

                # Build cell HTML
                attrs_str = " " + " ".join(cell_attrs) if cell_attrs else ""
                html_lines.append(f'    <{cell_tag}{attrs_str}>{cell_text}</{cell_tag}>')

                # Skip columns covered by colspan
                col_idx += colspan

            html_lines.append("  </tr>")

        html_lines.append("</table>")

        return "\n".join(html_lines) + "\n\n"

    def _extract_image_as_base64(self, image_data: bytes, image_format: str = 'png') -> str:
        """Convert image bytes to base64 string"""
        return base64.b64encode(image_data).decode('utf-8')

    def _ocr_image(self, image_bytes: bytes) -> str:
        """Use OCR to convert image to text description"""
        if not image_bytes:
            return "No image data"

        if not self.ocr:
            return "OCR not available"

        try:
            # Use the configured OCR instance with language configuration if available
            kwargs = {}
            if hasattr(self.ocr, 'config') and self.ocr.config and self.ocr.config.language:
                kwargs['language'] = self.ocr.config.language

            result = self.ocr.process_image(image_bytes, **kwargs)
            if hasattr(result, 'text'):
                return result.text
            else:
                return str(result)
        except Exception as e:
            logger.error(f"OCR failed: {e}")
            return "OCR failed"

    def _classify_text_type(self, text: str, style_name: str) -> str:
        """Classify text type based on style and content"""
        if not text:
            return "text:normal"

        # Check style-based classification first
        if style_name:
            style_lower = style_name.lower()

            # Check for title/heading styles
            if 'title' in style_lower:
                return "text:title"
            elif 'heading' in style_lower:
                # Extract heading level if possible
                if '1' in style_name:
                    return "text:title"
                else:
                    return "text:section"
            elif 'subtitle' in style_lower:
                return "text:section"
            elif 'caption' in style_lower:
                return "text:caption"
            elif any(x in style_lower for x in ['list', 'bullet']):
                return "text:list"

        # Content-based classification as fallback
        text_lower = text.lower()

        # Check for list patterns
        if re.match(r'^[\u2022•\-\*\d]+[\.\)]\s+', text):
            return "text:list"

        # Check for caption patterns
        caption_patterns = [
            r'^(Figure|Fig\.?|Table|Tbl\.?|Chart|Graph|Image|Plate|Scheme)\s*\d*[\.:)]?',
            r'^(Source|Note|Notes)[\.:)]',
        ]

        for pattern in caption_patterns:
            if re.match(pattern, text, re.IGNORECASE):
                return "text:caption"

        # Default to normal text
        return "text:normal"

    def _batch_ocr_images(self, images_info: List[Dict[str, Any]]) -> Dict[str, str]:
        """Process multiple images with OCR in a single batch call
        
        Args:
            images_info: List of dictionaries containing:
                - 'id': Unique identifier for the image
                - 'data': Image bytes
                
        Returns:
            Dictionary mapping image IDs to OCR text results
        """
        if not images_info:
            return {}

        if not self.ocr:
            logger.warning("No OCR instance available for batch processing")
            return {}

        try:
            # Prepare image data for batch processing
            image_data_list = [info['data'] for info in images_info]
            image_ids = [info['id'] for info in images_info]

            language_info = getattr(self.ocr.config, 'language', 'auto') if hasattr(self.ocr,
                                                                                    'config') and self.ocr.config else 'auto'
            logger.info(f"Processing {len(image_data_list)} images with configured OCR (language: {language_info})...")

            # Use the configured OCR instance for batch processing
            # Pass language configuration if available
            kwargs = {}
            if hasattr(self.ocr, 'config') and self.ocr.config and self.ocr.config.language:
                kwargs['language'] = self.ocr.config.language
                logger.info(f"🌍 Passing language configuration to OCR: {self.ocr.config.language}")

            # Always use batch processing
            ocr_results = self.ocr.batch_process_images(image_data_list, **kwargs)

            # Map results back using both hash and ID for duplicate handling
            # This ensures compatibility with individual lookup methods while preserving duplicates
            results_map = {}
            id_to_result = {}  # Additional map for ID-based lookup
            
            for image_info, ocr_result in zip(images_info, ocr_results):
                # Use hash of image data as primary key - this matches individual lookup
                img_hash = hash(image_info['data'])
                
                # Store the OCR result
                ocr_text = ocr_result.text if hasattr(ocr_result, 'text') else str(ocr_result)
                
                # Store by hash (for compatibility)
                results_map[img_hash] = ocr_text
                
                # Also store by ID (for handling duplicates)
                id_to_result[image_info['id']] = ocr_text

            # For XLSX fallback images, also store with special keys for duplicate handling
            for image_info, ocr_result in zip(images_info, ocr_results):
                if image_info.get('location', {}).get('source') == 'zip_fallback':
                    img_idx = image_info['location']['img_idx']
                    fallback_key = ('_fallback_ocr', img_idx)
                    ocr_text = ocr_result.text if hasattr(ocr_result, 'text') else str(ocr_result)
                    results_map[fallback_key] = ocr_text

            logger.info(f"Successfully processed {len(ocr_results)} images with OCR")
            return results_map

        except Exception as e:
            logger.error(f"Batch OCR processing failed: {e}")
            return {}

    def _collect_all_images(self) -> List[Dict[str, Any]]:
        """Collect all images from the document for batch processing
        
        This method should be overridden by subclasses
        
        Returns:
            List of dictionaries containing:
                - 'id': Unique identifier for the image
                - 'data': Image bytes
                - 'location': Location info (page/slide/sheet number, etc.)
        """
        raise NotImplementedError("Subclasses must implement _collect_all_images")


class DocxLoader(BaseOfficeLoader):
    """Loader for DOCX (Word) documents"""

    def __init__(self, file_path: Union[str, Path], ocr=None, table_style: Union[str, TableStyle] = None):
        super().__init__(file_path, ocr, table_style)
        self._open_document()

    def _open_document(self):
        """Open DOCX document with error handling and configuration logging"""
        try:
            self.doc = Document(self.file_path)

            # Log DOCX configuration
            logger.info("=" * 60)
            logger.info(f"DOCX Configuration for: {self.file_path.name}")
            logger.info("=" * 60)
            logger.info(f"File path: {self.file_path}")
            logger.info(f"File size: {self.file_path.stat().st_size / (1024 * 1024):.2f} MB")

            # Count paragraphs and tables
            para_count = len(list(self.doc.paragraphs))
            table_count = len(self.doc.tables)
            logger.info(f"Total paragraphs: {para_count}")
            logger.info(f"Total tables: {table_count}")

            # Count images (approximate - includes inline shapes)
            inline_shape_count = len(self.doc.inline_shapes) if hasattr(self.doc, 'inline_shapes') else 0
            logger.info(f"Inline shapes (includes images): {inline_shape_count}")

            # Count sections
            section_count = len(self.doc.sections)
            logger.info(f"Total sections: {section_count}")

            # Document properties
            core_props = self.doc.core_properties
            if core_props:
                logger.info("Document Properties:")
                if core_props.title:
                    logger.info(f"  Title: {core_props.title}")
                if core_props.author:
                    logger.info(f"  Author: {core_props.author}")
                if core_props.subject:
                    logger.info(f"  Subject: {core_props.subject}")
                if core_props.created:
                    logger.info(f"  Created: {core_props.created}")
                if core_props.modified:
                    logger.info(f"  Modified: {core_props.modified}")
                if core_props.last_modified_by:
                    logger.info(f"  Last modified by: {core_props.last_modified_by}")

            logger.info("=" * 60)

        except Exception as e:
            logger.error(f"Failed to open DOCX: {e}")
            raise

    def convert_to_json(self,
                        extract_images: bool = True,
                        ocr_images: bool = False,
                        show_progress: bool = True) -> Dict[str, Any]:
        """Convert DOCX to JSON format"""
        result = {
            "filename": self.file_path.name,
            "pages": 1,  # DOCX doesn't have fixed pages
            "content": []
        }

        if show_progress:
            logging.info(f"Processing DOCX: {self.file_path.name}")

        # Batch OCR processing if requested
        ocr_results_map = {}
        processed_image_hashes = set()  # Track processed images to avoid duplicates
        if extract_images and ocr_images:
            if show_progress:
                logger.info("Collecting all images for batch OCR processing...")

            all_images_info = self._collect_all_images()

            if all_images_info:
                if show_progress:
                    logger.info(f"Processing {len(all_images_info)} images with batch OCR...")

                # Prepare batch for OCR
                ocr_batch = []
                image_hashes = []  # Store hashes to map results back

                for info in all_images_info:
                    base64_data = base64.b64encode(info['data']).decode('utf-8')
                    ocr_batch.append({"image_data": base64_data})
                    # Use hash of image data as key
                    image_hashes.append(hash(info['data']))

                try:
                    # Use the configured OCR instance from BaseOfficeLoader
                    ocr_results_map = self._batch_ocr_images(all_images_info)

                    if show_progress:
                        logger.info(f"Successfully processed {len(ocr_results_map)} images with OCR")

                except Exception as e:
                    logger.error(f"Batch OCR processing failed: {e}")
                    ocr_images = False  # Fall back to base64 extraction

        # Process document body
        for element in self._iter_block_items():
            if isinstance(element, Paragraph):
                self._process_paragraph(element, result["content"], extract_images, ocr_images, ocr_results_map, processed_image_hashes)
            elif isinstance(element, Table):
                table_md = self._convert_table_to_markdown(element, extract_images, ocr_images, ocr_results_map)
                if table_md:
                    result["content"].append({
                        "type": "table",
                        "content": table_md
                    })

                # Note: Images are now handled within the table cells, no need to extract separately

        # Also check for inline shapes at document level
        # NOTE: This is now disabled to avoid duplicate OCR results
        # Images are already processed via paragraphs and tables
        # if extract_images and hasattr(self.doc, 'inline_shapes'):
        #     for inline_shape in self.doc.inline_shapes:
        #         if hasattr(inline_shape, '_inline'):
        #             image_content = self._extract_inline_shape_image(inline_shape, ocr_images, ocr_results_map)
        #             if image_content:
        #                 result["content"].append(image_content)

        # Extract images from headers and footers
        if extract_images:
            try:
                for section_idx, section in enumerate(self.doc.sections):
                    # Process header
                    if hasattr(section, 'header'):
                        header = section.header
                        for para in header.paragraphs:
                            self._process_paragraph(para, result["content"], extract_images, ocr_images,
                                                    ocr_results_map, processed_image_hashes)

                    # Process footer
                    if hasattr(section, 'footer'):
                        footer = section.footer
                        for para in footer.paragraphs:
                            self._process_paragraph(para, result["content"], extract_images, ocr_images,
                                                    ocr_results_map, processed_image_hashes)

            except Exception as e:
                logger.warning(f"Failed to process headers/footers: {e}")

        return result

    def _iter_block_items(self):
        """Yield each paragraph and table in document order"""
        parent = self.doc.element.body
        for child in parent.iterchildren():
            if isinstance(child, docx.oxml.text.paragraph.CT_P):
                yield Paragraph(child, self.doc)
            elif isinstance(child, docx.oxml.table.CT_Tbl):
                yield Table(child, self.doc)

    def _process_paragraph(self, paragraph: Paragraph, content: List[Dict], extract_images: bool,
                           ocr_images: bool = False, ocr_results_map: Dict[str, str] = {},
                           processed_image_hashes: set = None):
        """Process a paragraph and extract text and images"""
        if processed_image_hashes is None:
            processed_image_hashes = set()
            
        # Extract images from runs first
        if extract_images:
            for run in paragraph.runs:
                image_result = self._extract_run_images(run, ocr_images, ocr_results_map, processed_image_hashes)
                if image_result:
                    if isinstance(image_result, list):
                        content.extend(image_result)
                    else:
                        content.append(image_result)

        # Then extract text
        text = paragraph.text.strip()
        if text:
            text_type = self._classify_text_type(text, paragraph.style.name if paragraph.style else "")
            content.append({
                "type": text_type,
                "content": text
            })

    def _extract_run_images(self, run, ocr_images: bool = False, ocr_results_map: Dict[str, str] = {},
                           processed_image_hashes: set = None) -> Optional[Union[Dict[str, str], List[Dict[str, str]]]]:
        """Extract all images from a run.

        Returns a single dict if exactly one image is found, a list of dicts if multiple
        images are found in the run, or None if no images are present.
        """
        if processed_image_hashes is None:
            processed_image_hashes = set()
            
        try:
            # Access the underlying XML element
            r_element = run._element

            # Look for drawing elements in the run
            found_items: List[Dict[str, str]] = []
            for child in r_element:
                # Check if this is a w:drawing element
                if child.tag.endswith('}drawing'):
                    # Look for inline or anchored shapes within the drawing
                    for drawing_child in child:
                        if drawing_child.tag.endswith('}inline'):
                            # Found an inline shape, extract the image
                            image_data = self._extract_image_from_inline(
                                drawing_child, ocr_images, ocr_results_map, processed_image_hashes
                            )
                            if image_data:
                                found_items.append(image_data)
                        elif drawing_child.tag.endswith('}anchor'):
                            # Found an anchored/floating shape, extract the image
                            image_data = self._extract_image_from_anchor(
                                drawing_child, ocr_images, ocr_results_map, processed_image_hashes
                            )
                            if image_data:
                                found_items.append(image_data)
            if found_items:
                return found_items if len(found_items) > 1 else found_items[0]

        except Exception as e:
            logging.warning(f"Failed to extract images from run: {e}")

        return None

    def _extract_image_from_inline(self, inline_element, ocr_images: bool = False,
                                   ocr_results_map: Dict[str, str] = {}, 
                                   processed_image_hashes: set = None) -> Optional[Dict[str, str]]:
        """Extract image from an inline element"""
        if processed_image_hashes is None:
            processed_image_hashes = set()
            
        try:
            # Navigate through the inline shape structure to find the blip
            for child in inline_element:
                if child.tag.endswith('}graphic'):
                    for graphic_child in child:
                        if graphic_child.tag.endswith('}graphicData'):
                            for data_child in graphic_child:
                                if data_child.tag.endswith('}pic'):
                                    # Found picture element
                                    for pic_child in data_child:
                                        if pic_child.tag.endswith('}blipFill'):
                                            for blip_child in pic_child:
                                                if blip_child.tag.endswith('}blip'):
                                                    # Get the embed relationship ID
                                                    embed_attr = None
                                                    for attr_name, attr_value in blip_child.attrib.items():
                                                        if attr_name.endswith('}embed'):
                                                            embed_attr = attr_value
                                                            break

                                                    if embed_attr:
                                                        # Get image using relationship ID
                                                        return self._get_image_by_rid(embed_attr, ocr_images,
                                                                                      ocr_results_map, processed_image_hashes)
        except Exception as e:
            logging.warning(f"Failed to extract image from inline element: {e}")

        return None

    def _extract_image_from_anchor(self, anchor_element, ocr_images: bool = False,
                                   ocr_results_map: Dict[str, str] = {}, 
                                   processed_image_hashes: set = None) -> Optional[Dict[str, str]]:
        """Extract image from an anchor element and return formatted result"""
        if processed_image_hashes is None:
            processed_image_hashes = set()
            
        try:
            # Navigate through the anchor element structure to find the blip
            for child in anchor_element:
                if child.tag.endswith('}graphic'):
                    for graphic_child in child:
                        if graphic_child.tag.endswith('}graphicData'):
                            for data_child in graphic_child:
                                if data_child.tag.endswith('}pic'):
                                    # Found picture element
                                    for pic_child in data_child:
                                        if pic_child.tag.endswith('}blipFill'):
                                            for blip_child in pic_child:
                                                if blip_child.tag.endswith('}blip'):
                                                    # Get the embed relationship ID
                                                    embed_attr = None
                                                    for attr_name, attr_value in blip_child.attrib.items():
                                                        if attr_name.endswith('}embed'):
                                                            embed_attr = attr_value
                                                            break

                                                    if embed_attr:
                                                        # Get image using relationship ID
                                                        return self._get_image_by_rid(embed_attr, ocr_images,
                                                                                      ocr_results_map, processed_image_hashes)
        except Exception as e:
            logging.warning(f"Failed to extract image from anchor element: {e}")

        return None

    def _get_image_by_rid(self, rid: str, ocr_images: bool = False, ocr_results_map: Dict[str, str] = {},
                          processed_image_hashes: set = None) -> Optional[Dict[str, str]]:
        """Get image data using relationship ID"""
        if processed_image_hashes is None:
            processed_image_hashes = set()
            
        try:
            # Get the image part using the relationship ID
            image_part = self.doc.part.related_parts.get(rid)
            if image_part:
                image_bytes = image_part.blob
                
                # Check if this image has already been processed
                img_hash = hash(image_bytes)
                if img_hash in processed_image_hashes:
                    return None  # Skip already processed images
                
                # Mark this image as processed
                processed_image_hashes.add(img_hash)

                if ocr_images:
                    # Use image content hash to find OCR result (already calculated above)

                    if img_hash in ocr_results_map:
                        ocr_text = ocr_results_map[img_hash]
                        return {
                            "type": "text:image_description",
                            "content": f"<image_ocr_result>{ocr_text}</image_ocr_result>"
                        }
                    else:
                        # Fallback to individual OCR if not in batch results
                        logger.warning(f"OCR result not found for image with rid {rid}, using fallback OCR")
                        ocr_text = self._ocr_image(image_bytes)
                        if ocr_text:
                            return {
                                "type": "text:image_description",
                                "content": f"<image_ocr_result>{ocr_text}</image_ocr_result>"
                            }
                else:
                    # Return base64 encoded image
                    base64_image = base64.b64encode(image_bytes).decode('utf-8')
                    return {
                        "type": "image",
                        "content": base64_image
                    }
        except Exception as e:
            logging.warning(f"Failed to get image by rId {rid}: {e}")

        return None

    def _extract_inline_shape_image(self, inline_shape, ocr_images: bool = False,
                                    ocr_results_map: Dict[str, str] = {}) -> Optional[Dict[str, str]]:
        """Extract image from an InlineShape object"""
        try:
            # Get the inline element
            inline = inline_shape._inline

            # Use the same extraction method
            return self._extract_image_from_inline(inline, ocr_images, ocr_results_map)

        except Exception as e:
            logging.warning(f"Failed to extract image from inline shape: {e}")

        return None

    def _collect_all_images(self) -> List[Dict[str, Any]]:
        """Collect all images from DOCX document for batch processing"""
        images_info = []
        image_counter = 0
        seen_images = set()  # Track unique images to avoid duplicates

        # Helper function to add unique images
        def add_unique_image(image_data: bytes, location_info: Dict) -> None:
            nonlocal image_counter
            if image_data:
                # Use hash to identify unique images
                img_hash = hash(image_data)
                if img_hash not in seen_images:
                    seen_images.add(img_hash)
                    image_counter += 1
                    images_info.append({
                        'id': f'docx_image_{image_counter}',
                        'data': image_data,
                        'location': location_info
                    })

        # 1. Collect from paragraphs (inline images in runs)
        for element in self._iter_block_items():
            if isinstance(element, Paragraph):
                for run in element.runs:
                    # Check for images in runs
                    r_element = run._element
                    for child in r_element:
                        if child.tag.endswith('}drawing'):
                            for drawing_child in child:
                                if drawing_child.tag.endswith('}inline'):
                                    image_data = self._extract_image_data_from_inline(drawing_child)
                                    add_unique_image(image_data, {'type': 'paragraph_inline'})
                                elif drawing_child.tag.endswith('}anchor'):
                                    # Handle floating/anchored images
                                    image_data = self._extract_image_data_from_anchor(drawing_child)
                                    add_unique_image(image_data, {'type': 'paragraph_anchor'})

        # 2. Collect from document inline shapes
        if hasattr(self.doc, 'inline_shapes'):
            for idx, inline_shape in enumerate(self.doc.inline_shapes):
                if hasattr(inline_shape, '_inline'):
                    image_data = self._extract_image_data_from_inline(inline_shape._inline)
                    add_unique_image(image_data, {'type': 'document_inline_shape', 'index': idx})

        # 3. Collect from headers and footers
        try:
            # Check all sections
            for section_idx, section in enumerate(self.doc.sections):
                # Headers
                if hasattr(section, 'header'):
                    header = section.header
                    # Check paragraphs in header
                    for para in header.paragraphs:
                        for run in para.runs:
                            r_element = run._element
                            for child in r_element:
                                if child.tag.endswith('}drawing'):
                                    for drawing_child in child:
                                        if drawing_child.tag.endswith('}inline'):
                                            image_data = self._extract_image_data_from_inline(drawing_child)
                                            add_unique_image(image_data, {'type': 'header', 'section': section_idx})
                                        elif drawing_child.tag.endswith('}anchor'):
                                            image_data = self._extract_image_data_from_anchor(drawing_child)
                                            add_unique_image(image_data,
                                                             {'type': 'header_anchor', 'section': section_idx})

                # Footers
                if hasattr(section, 'footer'):
                    footer = section.footer
                    # Check paragraphs in footer
                    for para in footer.paragraphs:
                        for run in para.runs:
                            r_element = run._element
                            for child in r_element:
                                if child.tag.endswith('}drawing'):
                                    for drawing_child in child:
                                        if drawing_child.tag.endswith('}inline'):
                                            image_data = self._extract_image_data_from_inline(drawing_child)
                                            add_unique_image(image_data, {'type': 'footer', 'section': section_idx})
                                        elif drawing_child.tag.endswith('}anchor'):
                                            image_data = self._extract_image_data_from_anchor(drawing_child)
                                            add_unique_image(image_data,
                                                             {'type': 'footer_anchor', 'section': section_idx})
        except Exception as e:
            logger.warning(f"Failed to extract images from headers/footers: {e}")

        # 4. Check for images in shapes (text boxes, etc.)
        try:
            # Access the document's body element
            body = self.doc.element.body

            # Look for AlternateContent elements which often contain shapes
            for elem in body.iter():
                if elem.tag.endswith('}AlternateContent'):
                    # Check Choice elements for modern format
                    for choice in elem:
                        if choice.tag.endswith('}Choice'):
                            for child in choice:
                                if child.tag.endswith('}drawing'):
                                    for drawing_child in child:
                                        if drawing_child.tag.endswith('}inline'):
                                            image_data = self._extract_image_data_from_inline(drawing_child)
                                            add_unique_image(image_data, {'type': 'shape_inline'})
                                        elif drawing_child.tag.endswith('}anchor'):
                                            image_data = self._extract_image_data_from_anchor(drawing_child)
                                            add_unique_image(image_data, {'type': 'shape_anchor'})
        except Exception as e:
            logger.warning(f"Failed to extract images from shapes: {e}")

        # 5. Check for images in tables
        try:
            for table_idx, element in enumerate(self._iter_block_items()):
                if isinstance(element, Table):
                    for row_idx, row in enumerate(element.rows):
                        for cell_idx, cell in enumerate(row.cells):
                            for para in cell.paragraphs:
                                for run in para.runs:
                                    r_element = run._element
                                    for child in r_element:
                                        if child.tag.endswith('}drawing'):
                                            for drawing_child in child:
                                                if drawing_child.tag.endswith('}inline'):
                                                    image_data = self._extract_image_data_from_inline(drawing_child)
                                                    add_unique_image(image_data, {
                                                        'type': 'table_cell',
                                                        'table': table_idx,
                                                        'row': row_idx,
                                                        'cell': cell_idx
                                                    })
                                                elif drawing_child.tag.endswith('}anchor'):
                                                    image_data = self._extract_image_data_from_anchor(drawing_child)
                                                    add_unique_image(image_data, {
                                                        'type': 'table_cell_anchor',
                                                        'table': table_idx,
                                                        'row': row_idx,
                                                        'cell': cell_idx
                                                    })
        except Exception as e:
            logger.warning(f"Failed to extract images from tables: {e}")

        # 6. Try to get all relationships and check for image parts
        try:
            # Get all relationships from document part
            for rel_id, rel in self.doc.part.rels.items():
                if "image" in rel.reltype:
                    try:
                        image_part = rel.target_part
                        if hasattr(image_part, 'blob'):
                            add_unique_image(image_part.blob, {'type': 'relationship', 'rel_id': rel_id})
                    except:
                        pass
        except Exception as e:
            logger.warning(f"Failed to extract images from relationships: {e}")

        logger.info(f"Collected {len(images_info)} unique images from DOCX")
        return images_info

    def _extract_image_data_from_inline(self, inline_element) -> Optional[bytes]:
        """Extract raw image data from an inline element"""
        try:
            # Navigate through the inline shape structure to find the blip
            for child in inline_element:
                if child.tag.endswith('}graphic'):
                    for graphic_child in child:
                        if graphic_child.tag.endswith('}graphicData'):
                            for data_child in graphic_child:
                                if data_child.tag.endswith('}pic'):
                                    # Found picture element
                                    for pic_child in data_child:
                                        if pic_child.tag.endswith('}blipFill'):
                                            for blip_child in pic_child:
                                                if blip_child.tag.endswith('}blip'):
                                                    # Get the embed relationship ID
                                                    embed_attr = None
                                                    for attr_name, attr_value in blip_child.attrib.items():
                                                        if attr_name.endswith('}embed'):
                                                            embed_attr = attr_value
                                                            break

                                                    if embed_attr:
                                                        # Get image data using relationship ID
                                                        image_part = self.doc.part.related_parts.get(embed_attr)
                                                        if image_part:
                                                            return image_part.blob
        except Exception as e:
            logging.warning(f"Failed to extract image data from inline element: {e}")

        return None

    def _extract_image_data_from_anchor(self, anchor_element) -> Optional[bytes]:
        """Extract raw image data from an anchor element"""
        try:
            # Navigate through the anchor element structure to find the blip
            for child in anchor_element:
                if child.tag.endswith('}graphic'):
                    for graphic_child in child:
                        if graphic_child.tag.endswith('}graphicData'):
                            for data_child in graphic_child:
                                if data_child.tag.endswith('}pic'):
                                    # Found picture element
                                    for pic_child in data_child:
                                        if pic_child.tag.endswith('}blipFill'):
                                            for blip_child in pic_child:
                                                if blip_child.tag.endswith('}blip'):
                                                    # Get the embed relationship ID
                                                    embed_attr = None
                                                    for attr_name, attr_value in blip_child.attrib.items():
                                                        if attr_name.endswith('}embed'):
                                                            embed_attr = attr_value
                                                            break

                                                    if embed_attr:
                                                        # Get image data using relationship ID
                                                        image_part = self.doc.part.related_parts.get(embed_attr)
                                                        if image_part:
                                                            return image_part.blob
        except Exception as e:
            logging.warning(f"Failed to extract image data from anchor element: {e}")

        return None

    def _extract_images_from_paragraph(self, paragraph: Paragraph, content: List[Dict], extract_images: bool,
                                       ocr_images: bool = False, ocr_results_map: Dict[str, str] = {}):
        """Extract only images from a paragraph (used for table cells to avoid duplicate text)"""
        if extract_images:
            for run in paragraph.runs:
                image_content = self._extract_run_images(run, ocr_images, ocr_results_map)
                if image_content:
                    content.append(image_content)

    def _extract_cell_content_with_images(self, cell, extract_images: bool = True,
                                          ocr_images: bool = False, ocr_results_map: Dict[str, str] = {}) -> str:
        """Extract complete cell content including text and images (with OCR if enabled)
        
        Args:
            cell: DOCX table cell object
            extract_images: Whether to extract images
            ocr_images: Whether to use OCR on images
            ocr_results_map: Pre-computed OCR results map
            
        Returns:
            Combined cell content as string
        """
        cell_parts = []

        # Process each paragraph in the cell
        for para in cell.paragraphs:
            para_parts = []

            # Process each run in the paragraph
            for run in para.runs:
                # First check for images in the run
                if extract_images:
                    # Check for drawing elements
                    r_element = run._element
                    for child in r_element:
                        if child.tag.endswith('}drawing'):
                            for drawing_child in child:
                                image_bytes = None
                                if drawing_child.tag.endswith('}inline'):
                                    image_bytes = self._extract_image_data_from_inline(drawing_child)
                                elif drawing_child.tag.endswith('}anchor'):
                                    image_bytes = self._extract_image_data_from_anchor(drawing_child)

                                if image_bytes:
                                    if ocr_images:
                                        # Use OCR to get text
                                        img_hash = hash(image_bytes)
                                        if img_hash in ocr_results_map:
                                            ocr_text = ocr_results_map[img_hash]
                                            para_parts.append(f"[Image: {ocr_text}]")
                                        else:
                                            # Fallback to individual OCR
                                            ocr_text = self._ocr_image(image_bytes)
                                            if ocr_text:
                                                para_parts.append(f"[Image: {ocr_text}]")
                                            else:
                                                para_parts.append("[Image]")
                                    else:
                                        # Just mark as image placeholder
                                        para_parts.append("[Image]")

                # Then add the text content
                if run.text:
                    para_parts.append(run.text)

            # Combine paragraph parts
            if para_parts:
                para_text = "".join(para_parts).strip()
                if para_text:
                    cell_parts.append(para_text)

        # Join all paragraphs with newlines
        return "\n".join(cell_parts)


class PptxLoader(BaseOfficeLoader):
    """Loader for PPTX (PowerPoint) documents"""

    def __init__(self, file_path: Union[str, Path], ocr=None, table_style: Union[str, TableStyle] = None):
        super().__init__(file_path, ocr, table_style)
        self._open_document()

    def _open_document(self):
        """Open PPTX document"""
        try:
            self.doc = Presentation(self.file_path)
            logger.info(f"Opened PPTX: {self.file_path.name}")
            logger.info(f"Total slides: {len(self.doc.slides)}")
        except Exception as e:
            logger.error(f"Failed to open PPTX: {e}")
            raise

    def convert_to_json(self,
                        extract_images: bool = True,
                        ocr_images: bool = False,
                        show_progress: bool = True) -> Dict[str, Any]:
        """Convert PPTX to JSON format"""
        document = {
            "filename": self.file_path.name,
            "pages": len(self.doc.slides),
            "content": []
        }

        # Batch OCR processing if requested
        ocr_results_map = {}
        if extract_images and ocr_images:
            if show_progress:
                logger.info("Collecting all images for batch OCR processing...")

            all_images_info = self._collect_all_images()

            if all_images_info:
                if show_progress:
                    logger.info(f"Processing {len(all_images_info)} images with batch OCR...")

                # Prepare batch for OCR
                ocr_batch = []
                image_hashes = []  # Store hashes to map results back

                for info in all_images_info:
                    base64_data = base64.b64encode(info['data']).decode('utf-8')
                    ocr_batch.append({"image_data": base64_data})
                    # Use hash of image data as key
                    image_hashes.append(hash(info['data']))

                try:
                    # Use the configured OCR instance from BaseOfficeLoader
                    ocr_results_map = self._batch_ocr_images(all_images_info)

                    if show_progress:
                        logger.info(f"Successfully processed {len(ocr_results_map)} images with OCR")

                except Exception as e:
                    logger.error(f"Batch OCR processing failed: {e}")
                    ocr_images = False  # Fall back to base64 extraction

        # Process each slide
        for slide_idx, slide in enumerate(self.doc.slides):
            if show_progress:
                logger.info(f"Processing slide {slide_idx + 1}/{len(self.doc.slides)}")

            slide_content = self._process_slide(slide, slide_idx + 1, extract_images, ocr_images, ocr_results_map)
            document["content"].extend(slide_content)

            # Extract notes from the slide if present
            notes_content = self._extract_slide_notes(slide, slide_idx + 1)
            if notes_content:
                document["content"].append(notes_content)

        return document

    def _collect_all_images(self) -> List[Dict[str, Any]]:
        """Collect all images from PPTX presentation for batch processing"""
        images_info = []
        image_counter = 0

        for slide_idx, slide in enumerate(self.doc.slides):
            # Check placeholders
            for placeholder in slide.placeholders:
                if hasattr(placeholder, 'image'):
                    try:
                        image_data = placeholder.image.blob
                        image_counter += 1
                        images_info.append({
                            'id': f'pptx_slide{slide_idx + 1}_placeholder_{image_counter}',
                            'data': image_data,
                            'location': {'slide': slide_idx + 1, 'type': 'placeholder'}
                        })
                    except:
                        pass

            # Check all shapes
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_data = shape.image.blob
                        image_counter += 1
                        images_info.append({
                            'id': f'pptx_slide{slide_idx + 1}_shape{shape_idx}_{image_counter}',
                            'data': image_data,
                            'location': {'slide': slide_idx + 1, 'type': 'shape', 'shape_idx': shape_idx}
                        })
                    except:
                        pass

                # Check grouped shapes
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    group_images = self._collect_images_from_group(shape, slide_idx + 1, image_counter)
                    images_info.extend(group_images)
                    image_counter += len(group_images)

        logger.info(f"Collected {len(images_info)} images from PPTX")
        return images_info

    def _collect_images_from_group(self, group_shape, slide_num: int, counter_start: int) -> List[Dict[str, Any]]:
        """Recursively collect images from grouped shapes"""
        images = []
        counter = counter_start

        try:
            for shape in group_shape.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # Recursive call for nested groups
                    nested_images = self._collect_images_from_group(shape, slide_num, counter)
                    images.extend(nested_images)
                    counter += len(nested_images)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_data = shape.image.blob
                        counter += 1
                        images.append({
                            'id': f'pptx_slide{slide_num}_group_{counter}',
                            'data': image_data,
                            'location': {'slide': slide_num, 'type': 'grouped'}
                        })
                    except:
                        pass
        except:
            pass

        return images

    def _process_slide(self, slide, slide_num: int, extract_images: bool, ocr_images: bool,
                       ocr_results_map: Dict[str, str] = {}) -> List[Dict[str, Any]]:
        """Process a single slide"""
        content_items = []

        logger.info(f"Processing slide {slide_num}...")

        # First extract from placeholders (most structured content)
        placeholder_content = self._extract_from_placeholders(slide, slide_num, extract_images, ocr_images,
                                                              ocr_results_map)
        content_items.extend(placeholder_content)
        logger.info(f"  Extracted {len(placeholder_content)} items from placeholders")

        # Then extract from regular shapes (including grouped shapes)
        shape_content = self._extract_from_shapes(slide, slide_num, extract_images, ocr_images, ocr_results_map)
        content_items.extend(shape_content)
        logger.info(f"  Extracted {len(shape_content)} items from shapes")

        # Extract text from slide master/layout (headers, footers, page numbers)
        try:
            # Check for text in slide layout that might not be in placeholders
            if hasattr(slide, 'slide_layout'):
                layout = slide.slide_layout
                layout_shape_count = len(layout.shapes) if hasattr(layout, 'shapes') else 0
                logger.debug(f"  Checking {layout_shape_count} shapes in slide layout")

                # Look for footer/header text in layout
                for shape in layout.shapes:
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text and len(text) < 100:  # Usually footers/headers are short
                            # Check if this text is already captured
                            text_exists = any(item.get('content', '') == text for item in content_items)
                            if not text_exists:
                                content_items.append({
                                    "type": "text:caption",
                                    "content": text,
                                    "page": slide_num,
                                    "_top": 1000,  # Put at bottom
                                    "_left": 0
                                })
                                logger.debug(f"    Added layout text: '{text}'")
        except Exception as e:
            logger.warning(f"Error extracting layout text: {e}")

        # Sort by position to maintain reading order
        content_items.sort(key=lambda x: (x.get("_top", 0), x.get("_left", 0)))

        # Remove internal position markers
        for item in content_items:
            item.pop("_top", None)
            item.pop("_left", None)

        logger.info(f"  Total items extracted from slide {slide_num}: {len(content_items)}")
        return content_items

    def _extract_from_placeholders(self, slide, slide_num: int, extract_images: bool, ocr_images: bool,
                                   ocr_results_map: Dict[str, str] = {}) -> List[
        Dict[str, Any]]:
        """Extract content from slide placeholders"""
        content_items = []

        try:
            # Import placeholder type enum
            from pptx.enum.shapes import PP_PLACEHOLDER

            placeholder_count = len(slide.placeholders) if hasattr(slide, 'placeholders') else 0
            logger.debug(f"Slide {slide_num}: Found {placeholder_count} placeholders")

            for idx, placeholder in enumerate(slide.placeholders):
                try:
                    # Get placeholder type for better classification
                    ph_type = placeholder.placeholder_format.type
                    logger.debug(f"  Placeholder {idx}: type={ph_type}")

                    # Handle text placeholders
                    if placeholder.has_text_frame:
                        text_content = self._extract_text_from_placeholder(placeholder, ph_type)
                        if text_content:
                            text_content["page"] = slide_num
                            text_content["_top"] = placeholder.top if hasattr(placeholder, 'top') else 0
                            text_content["_left"] = placeholder.left if hasattr(placeholder, 'left') else 0
                            content_items.append(text_content)
                            logger.debug(
                                f"    Extracted text: {text_content['type']} - {len(text_content['content'])} chars")

                    # Handle table placeholders
                    if hasattr(placeholder, 'has_table') and placeholder.has_table:
                        table_content = self._extract_table_from_shape(placeholder, slide_num)
                        if table_content:
                            table_content["_top"] = placeholder.top if hasattr(placeholder, 'top') else 0
                            table_content["_left"] = placeholder.left if hasattr(placeholder, 'left') else 0
                            content_items.append(table_content)
                            logger.debug(f"    Extracted table")

                    # Handle picture placeholders
                    if extract_images and ph_type == PP_PLACEHOLDER.PICTURE:
                        # Picture placeholders might be populated
                        if hasattr(placeholder, 'image'):
                            image_content = self._extract_image_from_placeholder(placeholder, slide_num, ocr_images,
                                                                                 ocr_results_map)
                            if image_content:
                                image_content["_top"] = placeholder.top if hasattr(placeholder, 'top') else 0
                                image_content["_left"] = placeholder.left if hasattr(placeholder, 'left') else 0
                                content_items.append(image_content)
                                logger.debug(f"    Extracted image")

                except Exception as e:
                    logger.warning(f"Failed to process placeholder {idx}: {e}")

        except ImportError as e:
            logger.warning(f"PP_PLACEHOLDER enum not available: {e}, using basic extraction")
            # Fallback: extract from all placeholders without type information
            try:
                for idx, placeholder in enumerate(slide.placeholders):
                    if placeholder.has_text_frame:
                        text = placeholder.text_frame.text.strip()
                        if text:
                            content_items.append({
                                "type": self._classify_text_type(text, ""),
                                "content": text,
                                "page": slide_num,
                                "_top": placeholder.top if hasattr(placeholder, 'top') else 0,
                                "_left": placeholder.left if hasattr(placeholder, 'left') else 0
                            })
            except Exception as e:
                logger.error(f"Fallback placeholder extraction failed: {e}")

        logger.debug(f"Slide {slide_num}: Extracted {len(content_items)} items from placeholders")
        return content_items

    def _extract_text_from_placeholder(self, placeholder, ph_type) -> Optional[Dict[str, Any]]:
        """Extract text from a placeholder with proper type classification"""
        text_parts = []

        for paragraph in placeholder.text_frame.paragraphs:
            # Don't strip yet - check raw text first
            para_text = paragraph.text
            # Only skip if truly empty or just whitespace
            if para_text and not para_text.isspace():
                # Now we can strip for storage
                text_parts.append(para_text.strip())

        if text_parts:
            full_text = "\n".join(text_parts)

            # Don't skip short text
            if not full_text:
                return None

            # Import placeholder types
            try:
                from pptx.enum.shapes import PP_PLACEHOLDER

                # Map placeholder types to our text types
                if ph_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                    text_type = "text:title"
                elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                    text_type = "text:section"
                elif ph_type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT,
                                 PP_PLACEHOLDER.VERTICAL_BODY, PP_PLACEHOLDER.VERTICAL_OBJECT]:
                    # Further classify body text
                    text_type = self._classify_text_type(full_text, "")
                elif ph_type in [PP_PLACEHOLDER.DATE, PP_PLACEHOLDER.FOOTER,
                                 PP_PLACEHOLDER.HEADER, PP_PLACEHOLDER.SLIDE_NUMBER]:
                    text_type = "text:caption"
                else:
                    text_type = "text:normal"

                # Log very short text from placeholders
                if len(full_text) <= 5:
                    logger.debug(f"    Found short text in placeholder: '{full_text}' (type: {ph_type})")

            except:
                # Fallback classification
                text_type = self._classify_text_type(full_text, "")

            return {
                "type": text_type,
                "content": full_text
            }

        return None

    def _extract_from_shapes(self, slide, slide_num: int, extract_images: bool, ocr_images: bool,
                             ocr_results_map: Dict[str, str] = {}) -> List[
        Dict[str, Any]]:
        """Extract content from regular shapes (non-placeholders)"""
        content_items = []

        shape_count = len(slide.shapes) if hasattr(slide, 'shapes') else 0
        logger.debug(f"Slide {slide_num}: Processing {shape_count} shapes")

        for idx, shape in enumerate(slide.shapes):
            try:
                # Skip placeholders as they're already processed
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    logger.debug(f"  Shape {idx}: Skipping (is placeholder)")
                    continue

                # Process grouped shapes recursively
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    logger.debug(f"  Shape {idx}: Processing group shape")
                    group_content = self._extract_from_group_shape(shape, slide_num, extract_images, ocr_images,
                                                                   ocr_results_map)
                    content_items.extend(group_content)

                # Tables - extract full table structure (skip text extraction to avoid duplicates)
                if shape.has_table:
                    table_content = self._extract_table_from_shape(shape, slide_num)
                    if table_content:
                        table_content["_top"] = shape.top if hasattr(shape, 'top') else 0
                        table_content["_left"] = shape.left if hasattr(shape, 'left') else 0
                        content_items.append(table_content)
                        logger.debug(f"  Shape {idx}: Extracted table")
                else:
                    # Extract ALL text from non-table shapes using comprehensive method
                    all_text_items = self._extract_all_text_from_shape(shape, slide_num)
                    for text_item in all_text_items:
                        text_item["page"] = slide_num
                        text_item["_top"] = shape.top if hasattr(shape, 'top') else 0
                        text_item["_left"] = shape.left if hasattr(shape, 'left') else 0
                        content_items.append(text_item)

                # Pictures
                if extract_images and hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_content = self._extract_image_from_shape(shape, slide_num, ocr_images, ocr_results_map)
                    if image_content:
                        image_content["_top"] = shape.top if hasattr(shape, 'top') else 0
                        image_content["_left"] = shape.left if hasattr(shape, 'left') else 0
                        content_items.append(image_content)
                        logger.debug(f"  Shape {idx}: Extracted image")

                # Log unknown shape types for debugging
                if hasattr(shape, 'shape_type'):
                    shape_type_name = shape.shape_type
                    if shape_type_name not in [MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.PICTURE,
                                               MSO_SHAPE_TYPE.PLACEHOLDER, MSO_SHAPE_TYPE.TEXT_BOX]:
                        logger.debug(
                            f"  Shape {idx}: Type={shape_type_name}, Name={shape.name if hasattr(shape, 'name') else 'unnamed'}")

                        # Try to extract text from any shape as a last resort
                        if hasattr(shape, 'text') and shape.text.strip():
                            content_items.append({
                                "type": "text:normal",
                                "content": shape.text.strip(),
                                "page": slide_num,
                                "_top": shape.top if hasattr(shape, 'top') else 0,
                                "_left": shape.left if hasattr(shape, 'left') else 0
                            })
                            logger.debug(f"    Extracted text from unknown shape type: {shape.text[:50]}...")

            except Exception as e:
                logger.warning(f"Failed to process shape {idx}: {e}")

        logger.debug(f"Slide {slide_num}: Extracted {len(content_items)} items from shapes")
        return content_items

    def _extract_from_group_shape(self, group_shape, slide_num: int, extract_images: bool, ocr_images: bool,
                                  ocr_results_map: Dict[str, str] = {}) -> List[
        Dict[str, Any]]:
        """Recursively extract content from grouped shapes"""
        content_items = []

        try:
            for shape in group_shape.shapes:
                # Recursively handle nested groups
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    nested_content = self._extract_from_group_shape(shape, slide_num, extract_images, ocr_images,
                                                                    ocr_results_map)
                    content_items.extend(nested_content)

                # Extract text
                elif shape.has_text_frame:
                    text_content = self._extract_text_from_shape(shape)
                    if text_content:
                        text_content["page"] = slide_num
                        text_content["_top"] = shape.top if hasattr(shape, 'top') else 0
                        text_content["_left"] = shape.left if hasattr(shape, 'left') else 0
                        content_items.append(text_content)

                # Extract tables
                elif shape.has_table:
                    table_content = self._extract_table_from_shape(shape, slide_num)
                    if table_content:
                        table_content["_top"] = shape.top if hasattr(shape, 'top') else 0
                        table_content["_left"] = shape.left if hasattr(shape, 'left') else 0
                        content_items.append(table_content)

                # Extract images
                elif extract_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_content = self._extract_image_from_shape(shape, slide_num, ocr_images, ocr_results_map)
                    if image_content:
                        image_content["_top"] = shape.top if hasattr(shape, 'top') else 0
                        image_content["_left"] = shape.left if hasattr(shape, 'left') else 0
                        content_items.append(image_content)

        except Exception as e:
            logger.warning(f"Failed to process group shape: {e}")

        return content_items

    def _extract_text_from_shape(self, shape) -> Optional[Dict[str, Any]]:
        """Extract text from a shape"""
        text_parts = []

        for paragraph in shape.text_frame.paragraphs:
            # Don't strip yet - check raw text first
            para_text = paragraph.text
            # Only skip if truly empty or just whitespace
            if para_text and not para_text.isspace():
                # Now we can strip for storage
                text_parts.append(para_text.strip())

        if text_parts:
            full_text = "\n".join(text_parts)

            # Don't skip short text - it might be important (like page numbers, labels, etc.)
            if not full_text:
                return None

            # Try to get more context for classification
            shape_name = shape.name.lower() if hasattr(shape, 'name') else ""

            # Check if it's likely a title based on shape name
            if 'title' in shape_name:
                text_type = "text:title"
            elif 'subtitle' in shape_name:
                text_type = "text:section"
            else:
                text_type = self._classify_text_type(full_text, "")

            # Log very short text for debugging
            if len(full_text) <= 5:
                logger.debug(f"    Found short text: '{full_text}' in shape: {shape_name}")

            return {
                "type": text_type,
                "content": full_text
            }

        return None

    def _extract_chart_info(self, shape, slide_num: int) -> Optional[Dict[str, Any]]:
        """Extract basic information from a chart"""
        try:
            chart = shape.chart
            if hasattr(chart, 'chart_title') and chart.has_title:
                title_text = chart.chart_title.text_frame.text.strip()
                if title_text:
                    return {
                        "type": "text:caption",
                        "content": f"Chart: {title_text}",
                        "page": slide_num
                    }
        except Exception as e:
            logger.warning(f"Failed to extract chart info: {e}")

        return None

    def _extract_slide_notes(self, slide, slide_num: int) -> Optional[Dict[str, Any]]:
        """Extract notes from a slide"""
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text = notes_slide.notes_text_frame.text.strip()

                if notes_text:
                    return {
                        "type": "text:normal",
                        "content": f"[Slide {slide_num} Notes]\n{notes_text}",
                        "page": slide_num
                    }
        except Exception as e:
            logger.warning(f"Failed to extract slide notes: {e}")

        return None

    def _extract_table_from_shape(self, shape, slide_num: int) -> Optional[Dict[str, Any]]:
        """Extract table from a shape with enhanced merged cell detection for PPTX"""
        table = shape.table

        # Get table dimensions
        num_rows = len(table.rows)
        num_cols = len(table.columns)

        logger.debug(f"Extracting table from slide {slide_num}: {num_rows}x{num_cols}")

        # Initialize table data with proper dimensions
        table_data = [['' for _ in range(num_cols)] for _ in range(num_rows)]
        merged_cells_info = []

        # Track which cells we've already processed as part of merges
        processed_cells = set()

        # Process each cell to detect merges using gridSpan and vMerge properties
        for row_idx in range(num_rows):
            for col_idx in range(num_cols):
                if (row_idx, col_idx) in processed_cells:
                    continue

                try:
                    # Get the cell at this position
                    cell = table.cell(row_idx, col_idx)
                    tc = cell._tc  # Get the table cell element

                    # Get span properties
                    gridSpan = tc.gridSpan  # Horizontal span (>1 means merged horizontally)
                    vMerge = tc.vMerge  # Vertical merge (True for continuation cells)

                    # Get cell text
                    cell_text = cell.text.strip() if hasattr(cell, 'text') else ''

                    # Check if this is a vertically merged cell continuation
                    if vMerge:
                        # This cell is part of a vertical merge from above
                        # Find the origin cell by looking upward
                        origin_row = row_idx - 1
                        while origin_row >= 0:
                            origin_cell = table.cell(origin_row, col_idx)
                            if not origin_cell._tc.vMerge:
                                # Found the origin of the vertical merge
                                break
                            origin_row -= 1

                        # Mark this cell as processed (it's empty in a vertical merge)
                        processed_cells.add((row_idx, col_idx))
                        logger.debug(f"Cell({row_idx},{col_idx}) is part of vertical merge from row {origin_row}")
                        continue

                    # Determine the span of this cell
                    colspan = gridSpan if gridSpan > 1 else 1
                    rowspan = 1

                    # If this cell has text and cells below are empty with vMerge=True, it's a vertical merge origin
                    if row_idx < num_rows - 1:
                        # Check cells below for vertical merge
                        check_row = row_idx + 1
                        while check_row < num_rows:
                            check_cell = table.cell(check_row, col_idx)
                            if check_cell._tc.vMerge:
                                rowspan += 1
                                check_row += 1
                            else:
                                break

                    # Place text in the current cell
                    table_data[row_idx][col_idx] = cell_text

                    # Mark all cells covered by this merge as processed
                    for r in range(row_idx, row_idx + rowspan):
                        for c in range(col_idx, col_idx + colspan):
                            processed_cells.add((r, c))
                            # Clear cells that are part of the merge (except origin)
                            if (r, c) != (row_idx, col_idx):
                                if r < num_rows and c < num_cols:
                                    table_data[r][c] = ''

                    # Record merge info if this is a merged cell
                    if rowspan > 1 or colspan > 1:
                        merged_cells_info.append({
                            'row': row_idx,
                            'col': col_idx,
                            'rowspan': rowspan,
                            'colspan': colspan
                        })
                        logger.debug(
                            f"Found merged cell at ({row_idx}, {col_idx}) with span {rowspan}x{colspan}, text: '{cell_text[:30]}...'")

                except Exception as e:
                    logger.warning(f"Error processing cell at ({row_idx}, {col_idx}): {e}")
                    table_data[row_idx][col_idx] = ''
                    processed_cells.add((row_idx, col_idx))

        # Log the extraction result
        logger.info(f"PPTX Table: Extracted {num_rows}x{num_cols} table with {len(merged_cells_info)} merged cells")

        # Debug: Log first few rows of extracted data
        if logger.isEnabledFor(logging.DEBUG):
            for i, row in enumerate(table_data[:3]):
                logger.debug(f"  Row {i}: {[cell[:20] + '...' if cell and len(cell) > 20 else cell for cell in row]}")

        # For PowerPoint tables, we'll always use HTML format for better structure preservation
        if table_data:
            # Create table info for HTML conversion
            table_info = {
                'is_complex': True,  # Always treat PowerPoint tables as complex
                'merged_cells': merged_cells_info,
                'row_count': num_rows,
                'col_count': num_cols,
                'cell_spans': {}
            }

            # Add span info
            for merge_info in merged_cells_info:
                key = (merge_info['row'], merge_info['col'])
                table_info['cell_spans'][key] = (merge_info['rowspan'], merge_info['colspan'])

            # Use HTML converter for PowerPoint tables
            html_table = self._convert_table_to_html(table_data, table_info)

            return {
                "type": "table",
                "content": html_table,
                "page": slide_num
            }

        return None

    def _extract_image_from_placeholder(self, placeholder, slide_num: int, ocr_images: bool,
                                        ocr_results_map: Dict[str, str] = {}) -> Optional[
        Dict[str, Any]]:
        """Extract image from a picture placeholder"""
        try:
            image = placeholder.image
            image_data = image.blob

            if ocr_images and self.ocr:
                # Use image content hash to find OCR result
                img_hash = hash(image_data)

                if img_hash in ocr_results_map:
                    ocr_text = ocr_results_map[img_hash]
                    return {
                        "type": "text:image_description",
                        "content": f"<image_ocr_result>{ocr_text}</image_ocr_result>",
                        "page": slide_num
                    }
                else:
                    # Fallback to individual OCR
                    logger.warning(
                        f"OCR result not found for placeholder image on slide {slide_num}, using fallback OCR")
                    ocr_text = self._ocr_image(image_data)
                    return {
                        "type": "text:image_description",
                        "content": f"<image_ocr_result>{ocr_text}</image_ocr_result>",
                        "page": slide_num
                    }
            else:
                base64_data = self._extract_image_as_base64(image_data)
                return {
                    "type": "image",
                    "content": base64_data,
                    "page": slide_num
                }
        except Exception as e:
            logger.warning(f"Failed to extract image from placeholder: {e}")
            return None

    def _extract_image_from_shape(self, shape, slide_num: int, ocr_images: bool,
                                  ocr_results_map: Dict[str, str] = {}) -> Optional[Dict[str, Any]]:
        """Extract image from a shape"""
        try:
            image = shape.image
            image_data = image.blob

            if ocr_images and self.ocr:
                # Use image content hash to find OCR result
                img_hash = hash(image_data)

                if img_hash in ocr_results_map:
                    ocr_text = ocr_results_map[img_hash]
                    return {
                        "type": "text:image_description",
                        "content": f"<image_ocr_result>{ocr_text}</image_ocr_result>",
                        "page": slide_num
                    }
                else:
                    # Fallback to individual OCR
                    logger.warning(f"OCR result not found for shape image on slide {slide_num}, using fallback OCR")
                    ocr_text = self._ocr_image(image_data)
                    return {
                        "type": "text:image_description",
                        "content": f"<image_ocr_result>{ocr_text}</image_ocr_result>",
                        "page": slide_num
                    }
            else:
                # Return base64 encoded image
                base64_data = self._extract_image_as_base64(image_data)
                return {
                    "type": "image",
                    "content": base64_data,
                    "page": slide_num
                }

        except Exception as e:
            logger.warning(f"Failed to extract image: {e}")
            return None

    def _extract_all_text_from_shape(self, shape, slide_num: int) -> List[Dict[str, Any]]:
        """Extract all possible text from any shape type"""
        text_items = []

        try:
            # 1. Regular text frame
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text_content = self._extract_text_from_shape(shape)
                if text_content:
                    text_items.append(text_content)

            # 2. Chart text elements (removed table cell extraction - tables handled separately)
            if hasattr(shape, 'has_chart') and shape.has_chart:
                chart = shape.chart

                # Chart title
                if hasattr(chart, 'has_title') and chart.has_title:
                    try:
                        title_text = chart.chart_title.text_frame.text.strip()
                        if title_text:
                            text_items.append({
                                "type": "text:caption",
                                "content": f"Chart: {title_text}"
                            })
                    except:
                        pass

                # Axis titles
                try:
                    if hasattr(chart, 'category_axis') and chart.category_axis.has_title:
                        axis_text = chart.category_axis.axis_title.text_frame.text.strip()
                        if axis_text:
                            text_items.append({
                                "type": "text:caption",
                                "content": f"X-axis: {axis_text}"
                            })
                except:
                    pass

                try:
                    if hasattr(chart, 'value_axis') and chart.value_axis.has_title:
                        axis_text = chart.value_axis.axis_title.text_frame.text.strip()
                        if axis_text:
                            text_items.append({
                                "type": "text:caption",
                                "content": f"Y-axis: {axis_text}"
                            })
                except:
                    pass

            # 3. SmartArt text (often in grouped shapes)
            # SmartArt is typically a group shape with specific properties
            if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Check if it might be SmartArt by looking at shape name
                shape_name = shape.name.lower() if hasattr(shape, 'name') else ""
                if 'diagram' in shape_name or 'smart' in shape_name:
                    logger.debug(f"  Possible SmartArt detected: {shape.name}")

            # 4. Shape alt text (often contains descriptions)
            if hasattr(shape, 'alt_text') and shape.alt_text:
                text_items.append({
                    "type": "text:caption",
                    "content": f"[Alt text]: {shape.alt_text}"
                })

            # 5. Connector text (lines with labels)
            if hasattr(shape, 'connector_type'):
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    connector_text = shape.text_frame.text.strip()
                    if connector_text:
                        text_items.append({
                            "type": "text:caption",
                            "content": connector_text
                        })

            # 6. OLE objects might have accessible text
            if hasattr(shape, 'ole_format'):
                logger.debug(f"  Found OLE object: {shape.name if hasattr(shape, 'name') else 'unnamed'}")
                # OLE objects are embedded files, limited text extraction

        except Exception as e:
            logger.warning(f"Error extracting text from shape: {e}")

        return text_items

    def _extract_master_layout_text(self, slide, slide_num: int) -> List[Dict[str, Any]]:
        """Extract text from slide master and layout that appears on the slide"""
        text_items = []

        try:
            # Get the slide layout
            slide_layout = slide.slide_layout

            # Extract text from layout placeholders that might not be filled in the slide
            for layout_placeholder in slide_layout.placeholders:
                try:
                    # Check if this placeholder exists in the slide
                    slide_has_placeholder = False
                    for slide_placeholder in slide.placeholders:
                        if slide_placeholder.placeholder_format.idx == layout_placeholder.placeholder_format.idx:
                            slide_has_placeholder = True
                            break

                    # If not in slide, check if layout has default text
                    if not slide_has_placeholder and layout_placeholder.has_text_frame:
                        layout_text = layout_placeholder.text_frame.text.strip()
                        if layout_text:
                            ph_type = layout_placeholder.placeholder_format.type
                            text_items.append({
                                "type": "text:caption",
                                "content": f"[Layout default]: {layout_text}",
                                "page": slide_num
                            })
                            logger.debug(f"  Found layout text: '{layout_text}'")

                except Exception as e:
                    logger.warning(f"Error extracting layout placeholder text: {e}")

            # Check slide master for footer/header text
            try:
                slide_master = slide_layout.slide_master
                for master_placeholder in slide_master.placeholders:
                    if master_placeholder.has_text_frame:
                        master_text = master_placeholder.text_frame.text.strip()
                        if master_text:
                            # Check if it's footer/header/date/slide number
                            ph_type = master_placeholder.placeholder_format.type
                            try:
                                from pptx.enum.shapes import PP_PLACEHOLDER
                                if ph_type in [PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.DATE,
                                               PP_PLACEHOLDER.SLIDE_NUMBER]:
                                    text_items.append({
                                        "type": "text:caption",
                                        "content": master_text,
                                        "page": slide_num
                                    })
                                    logger.debug(f"  Found master text: '{master_text}' (type: {ph_type})")
                            except:
                                pass

            except Exception as e:
                logger.warning(f"Error extracting master text: {e}")

        except Exception as e:
            logger.warning(f"Error accessing slide layout/master: {e}")

        return text_items


class XlsxLoader(BaseOfficeLoader):
    """Loader for XLSX (Excel) documents"""

    def __init__(self, file_path: Union[str, Path], ocr=None, table_style: Union[str, TableStyle] = None):
        super().__init__(file_path, ocr, table_style)
        self._open_document()

    def _open_document(self):
        """Open XLSX document"""
        try:
            self.doc = openpyxl.load_workbook(self.file_path, data_only=True)
            self.df_sheets = pd.read_excel(self.file_path, sheet_name=None)
            logger.info(f"Opened XLSX: {self.file_path.name}")
            logger.info(f"Total sheets: {len(self.doc.worksheets)}")
        except Exception as e:
            logger.error(f"Failed to open XLSX: {e}")
            raise

    def convert_to_json(self,
                        extract_images: bool = True,
                        ocr_images: bool = False,
                        show_progress: bool = True) -> Dict[str, Any]:
        """Convert XLSX to JSON format"""
        document = {
            "filename": self.file_path.name,
            "pages": len(self.doc.worksheets),
            "content": []
        }

        # Batch OCR processing if requested
        ocr_results_map = {}
        # Store image data to avoid calling _data() twice
        image_data_cache = {}  
        
        if extract_images and ocr_images:
            if show_progress:
                logger.info("Collecting all images for batch OCR processing...")

            all_images_info = self._collect_all_images()
            
            # Cache image data for reuse during sheet extraction
            # Build cache using sheet_name and img_idx from the location info
            for info in all_images_info:
                sheet_name = info['location']['sheet_name']
                img_idx = info['location']['img_idx']
                cache_key = (sheet_name, img_idx)
                image_data_cache[cache_key] = info['data']
                
                # Also cache fallback images with a special key
                if info['location'].get('source') == 'zip_fallback':
                    fallback_key = ('_fallback', img_idx)
                    image_data_cache[fallback_key] = info['data']

            if all_images_info:
                if show_progress:
                    logger.info(f"Processing {len(all_images_info)} images with batch OCR...")

                # Prepare batch for OCR
                ocr_batch = []
                image_hashes = []  # Store hashes to map results back

                for info in all_images_info:
                    base64_data = base64.b64encode(info['data']).decode('utf-8')
                    ocr_batch.append({"image_data": base64_data})
                    # Use hash of image data as key
                    image_hashes.append(hash(info['data']))

                try:
                    # Use the configured OCR instance from BaseOfficeLoader
                    ocr_results_map = self._batch_ocr_images(all_images_info)

                    if show_progress:
                        logger.info(f"Successfully processed {len(ocr_results_map)} images with OCR")

                except Exception as e:
                    logger.error(f"Batch OCR processing failed: {e}")
                    ocr_images = False  # Fall back to base64 extraction
        elif extract_images:
            # Even without OCR, collect and cache images to avoid double _data() calls
            all_images_info = self._collect_all_images()
            # Build cache using sheet_name and img_idx from the location info
            for info in all_images_info:
                sheet_name = info['location']['sheet_name']
                img_idx = info['location']['img_idx']
                cache_key = (sheet_name, img_idx)
                image_data_cache[cache_key] = info['data']
                
                # Also cache fallback images with a special key
                if info['location'].get('source') == 'zip_fallback':
                    fallback_key = ('_fallback', img_idx)
                    image_data_cache[fallback_key] = info['data']

        # Track which images have been embedded in tables to avoid duplicates
        embedded_image_hashes = set()
        
        # Process each worksheet
        for sheet_idx, sheet_name in enumerate(self.doc.sheetnames):
            if show_progress:
                logger.info(f"Processing sheet {sheet_idx + 1}/{len(self.doc.worksheets)}: {sheet_name}")

            sheet = self.doc[sheet_name]

            # Add sheet title
            document["content"].append({
                "type": "text:title",
                "content": f"Sheet: {sheet_name}",
                "page": sheet_idx + 1
            })

            # Extract table data with merged cell detection
            table_content = self._extract_sheet_as_table(sheet, sheet_name)
            
            # If we have OCR results and the table contains #VALUE!, embed OCR in the table
            if table_content and ocr_images and ocr_results_map and "#VALUE!" in table_content:
                # Get the first available OCR result for #VALUE! replacement
                for info in all_images_info:
                    if info['location']['sheet_name'] == sheet_name:
                        img_hash = hash(info['data'])
                        if img_hash in ocr_results_map:
                            ocr_text = ocr_results_map[img_hash]
                            
                            # Convert OCR result to structured single-line text for AI agents
                            import html
                            import re
                            
                            # Remove markdown formatting to convert to pure text
                            pure_text = ocr_text
                            
                            # Remove markdown bold/italic markers
                            pure_text = re.sub(r'\*\*([^*]+)\*\*', r'\1', pure_text)  # Remove bold **text**
                            pure_text = re.sub(r'\*([^*]+)\*', r'\1', pure_text)      # Remove italic *text*
                            pure_text = re.sub(r'__([^_]+)__', r'\1', pure_text)      # Remove bold __text__
                            pure_text = re.sub(r'_([^_]+)_', r'\1', pure_text)        # Remove italic _text_
                            
                            # Convert headers to structured sections
                            pure_text = re.sub(r'^#{1,6}\s+(.+)$', r'[\1]:', pure_text, flags=re.MULTILINE)
                            
                            # Remove code blocks markers but keep content
                            pure_text = re.sub(r'```[^\n]*\n', '', pure_text)
                            pure_text = re.sub(r'```', '', pure_text)
                            pure_text = re.sub(r'`([^`]+)`', r'\1', pure_text)  # Remove inline code markers
                            
                            # Convert markdown lists to structured format
                            pure_text = re.sub(r'^[\s]*[-*+]\s+(.+)$', r'• \1', pure_text, flags=re.MULTILINE)
                            pure_text = re.sub(r'^[\s]*\d+\.\s+(.+)$', r'• \1', pure_text, flags=re.MULTILINE)
                            
                            # Structure the text for AI understanding
                            # Split by common OCR sections
                            sections = []
                            
                            # Check for common patterns in OCR results
                            if "Extracted Text:" in pure_text or "[Extracted Text]:" in pure_text:
                                # Split into logical sections
                                lines = pure_text.split('\n')
                                current_section = []
                                for line in lines:
                                    line = line.strip()
                                    if line and (line.startswith('[') or 'Analysis:' in line or 'Summary:' in line or 'Text:' in line):
                                        if current_section:
                                            sections.append(' '.join(current_section))
                                        current_section = [line]
                                    elif line:
                                        current_section.append(line)
                                if current_section:
                                    sections.append(' '.join(current_section))
                                
                                # Join sections with clear separators
                                pure_text = ' || '.join(sections) if sections else pure_text
                            else:
                                # For simple text, just clean up line breaks
                                pure_text = re.sub(r'\n+', ' ', pure_text)
                            
                            # Clean up multiple spaces and format
                            pure_text = re.sub(r'\s+', ' ', pure_text)
                            pure_text = re.sub(r'\|\|\s+\|\|', '||', pure_text)  # Clean up empty sections
                            pure_text = pure_text.strip()
                            
                            # Escape HTML entities
                            pure_text_escaped = html.escape(pure_text)
                            
                            # Create a formatted div with structured OCR result
                            cell_content = f'''<div style="max-width: 100%; overflow-x: auto; background: #f9f9f9; padding: 8px; border: 1px solid #ddd; border-radius: 4px; font-family: 'Segoe UI', Arial, sans-serif; font-size: 12px; line-height: 1.4; color: #333;">
<strong>📷 OCR Analysis:</strong> <span style="font-family: 'Consolas', monospace; color: #0066cc;">{pure_text_escaped}</span>
</div>'''
                            
                            # Replace #VALUE! with the FULL OCR result
                            # Replace only the first placeholder occurrence on this sheet
                            table_content = table_content.replace(
                                "#VALUE!</td>",
                                f"{cell_content}</td>",
                                1
                            )
                            
                            # Mark that we've embedded this image's OCR
                            embedded_image_hashes.add(img_hash)
                            logger.info(f"Replaced #VALUE! with full OCR result in sheet {sheet_name}")
                            break
            
            if table_content:
                document["content"].append({
                    "type": "table",
                    "content": table_content,
                    "page": sheet_idx + 1
                })

            # Only extract images that weren't already embedded in table cells
            if extract_images:
                # Filter out images that were already embedded
                filtered_images = []
                images = self._extract_images_from_sheet(sheet, sheet_idx + 1, ocr_images, ocr_results_map, image_data_cache)
                
                for img in images:
                    # Check if this image's OCR was already embedded
                    img_content = img.get('content', '')
                    skip_image = False
                    
                    # Check if any embedded hash matches this image's content
                    for embedded_hash in embedded_image_hashes:
                        if embedded_hash in ocr_results_map:
                            embedded_text = ocr_results_map[embedded_hash]
                            # If this image's content matches an embedded one, skip it
                            if embedded_text in img_content:
                                skip_image = True
                                break
                    
                    if not skip_image:
                        filtered_images.append(img)
                
                if filtered_images:
                    document["content"].extend(filtered_images)

        return document

    def _extract_sheet_as_table(self, sheet, sheet_name: str) -> str:
        """Extract sheet data as table with merged cell detection"""
        # Get DataFrame for data
        if sheet_name in self.df_sheets:
            df = self.df_sheets[sheet_name]
            if df.empty:
                return ""
            
            # Drop completely empty columns from DataFrame
            df = df.dropna(axis=1, how='all')
        else:
            return ""

        # Check for merged cells in the worksheet
        merged_cells_ranges = sheet.merged_cells.ranges if hasattr(sheet, 'merged_cells') else []
        has_merged_cells = len(merged_cells_ranges) > 0

        if has_merged_cells:
            # Extract with merged cell information
            return self._extract_sheet_with_merged_cells(sheet, df)
        else:
            # Use standard DataFrame conversion
            return self._dataframe_to_markdown(df)

    def _extract_sheet_with_merged_cells(self, sheet, df: pd.DataFrame) -> str:
        """Extract sheet data handling merged cells"""
        # Get the actual dimensions
        max_row = sheet.max_row
        
        # Find actual max column with data across ALL rows (not Excel's 16384 limit)
        # This handles sheets with multiple tables that have different column counts
        actual_max_col = 1
        
        # Check ALL rows to find the true max column with data
        for row in sheet.iter_rows(min_row=1, max_row=max_row):
            for idx, cell in enumerate(row, 1):
                if cell.value is not None:
                    actual_max_col = max(actual_max_col, idx)
        
        # Also check merged cells, but exclude the buggy full-width merges (A:XFD)
        for merge_range in sheet.merged_cells.ranges:
            # Skip merged cells that span to Excel's max column (XFD = 16384)
            # These are typically separator rows, not real data
            if merge_range.max_col >= 16384:
                logger.debug(f"Skipping full-width merged cell: {merge_range}")
                continue
            
            # For normal merged cells, include their extent
            actual_max_col = max(actual_max_col, merge_range.max_col)
        
        max_col = actual_max_col
        logger.debug(f"Sheet dimensions: {max_row} rows x {max_col} columns (actual used)")

        # Build table data with merged cell info
        table_data = []
        merged_cells_info = []

        # Map merged cell ranges
        merge_map = {}
        for merge_range in sheet.merged_cells.ranges:
            # Get the top-left cell of the merge
            min_row = merge_range.min_row
            min_col = merge_range.min_col
            max_row_merge = merge_range.max_row
            # Limit merged cell columns to actual data bounds
            # (Handles buggy Excel files with merges to column XFD)
            max_col_merge = min(merge_range.max_col, max_col)

            # Get the value from the top-left cell
            cell_value = sheet.cell(row=min_row, column=min_col).value

            # Store merge info - only mark the origin cell
            merge_map[(min_row, min_col)] = {
                'min_row': min_row,
                'min_col': min_col,
                'max_row': max_row_merge,
                'max_col': max_col_merge,
                'value': cell_value,
                'rowspan': max_row_merge - min_row + 1,
                'colspan': max_col_merge - min_col + 1
            }

            # Mark other cells in the merge as "merged" (not origin)
            for r in range(min_row, max_row_merge + 1):
                for c in range(min_col, max_col_merge + 1):
                    if (r, c) != (min_row, min_col):
                        # Keep a backlink to the origin so we can map image anchors correctly
                        merge_map[(r, c)] = {'is_merged': True, 'origin': (min_row, min_col)}

        # Detect cells containing images (based on anchor positions) so we can embed OCR in-place
        image_origin_cells = set()
        try:
            from openpyxl.utils.cell import coordinate_from_string, column_index_from_string
        except Exception:
            coordinate_from_string = None
            column_index_from_string = None

        for _img in getattr(sheet, "_images", []):
            row = None
            col = None
            anchor = getattr(_img, "anchor", None)
            try:
                # Anchor can be a coordinate string like 'B3'
                if isinstance(anchor, str) and coordinate_from_string is not None:
                    col_letter, row_num = coordinate_from_string(anchor)
                    row = int(row_num)
                    col = int(column_index_from_string(col_letter))
                # Or an anchor object (OneCellAnchor/TwoCellAnchor) with 0-based indices
                elif hasattr(anchor, "_from") and hasattr(anchor._from, "row") and hasattr(anchor._from, "col"):
                    row = int(anchor._from.row) + 1
                    col = int(anchor._from.col) + 1
            except Exception:
                row = None
                col = None

            if row is None or col is None:
                continue

            # If the anchor falls inside a merged region, use that region's origin
            origin_row, origin_col = row, col
            mi = merge_map.get((row, col))
            if isinstance(mi, dict) and mi.get('is_merged') and 'origin' in mi:
                origin_row, origin_col = mi['origin']
            else:
                # Search origins in merge_map
                for _pos, _info in merge_map.items():
                    if 'min_row' in _info:
                        if _info['min_row'] <= row <= _info['max_row'] and _info['min_col'] <= col <= _info['max_col']:
                            origin_row, origin_col = _info['min_row'], _info['min_col']
                            break

            image_origin_cells.add((origin_row, origin_col))

        # Extract data row by row
        # Start from row 1 (Excel is 1-indexed)
        for row_idx in range(1, max_row + 1):
            row_data = []

            for col_idx in range(1, max_col + 1):
                # Check if this cell is part of a merge
                if (row_idx, col_idx) in merge_map:
                    merge_info = merge_map[(row_idx, col_idx)]

                    if 'is_merged' in merge_info:
                        # This cell is part of a merge but not the origin
                        row_data.append("")
                    else:
                        # This is the origin cell of a merge
                        value = str(merge_info['value']) if merge_info['value'] is not None else ""
                        # If an image is anchored to this merged region's origin, mark placeholder for OCR embedding
                        if (row_idx, col_idx) in image_origin_cells:
                            value = "#VALUE!" if not value else f"{value} #VALUE!"
                        row_data.append(value)

                        # Record merge info (0-indexed for our table)
                        if merge_info['rowspan'] > 1 or merge_info['colspan'] > 1:
                            merged_cells_info.append({
                                'row': len(table_data),  # Current row in table_data
                                'col': col_idx - 1,  # Convert to 0-indexed
                                'rowspan': merge_info['rowspan'],
                                'colspan': merge_info['colspan']
                            })
                else:
                    # Regular cell - get its value
                    cell_value = sheet.cell(row=row_idx, column=col_idx).value
                    value = str(cell_value) if cell_value is not None else ""
                    # If an image is anchored to this exact cell, mark placeholder for OCR embedding
                    if (row_idx, col_idx) in image_origin_cells:
                        value = "#VALUE!" if not value else f"{value} #VALUE!"
                    row_data.append(value)

            # Add all rows to table_data initially
            table_data.append(row_data)

        # Filter out columns that are completely empty across ALL rows
        # This preserves columns that have data in ANY table within the sheet
        if table_data and len(table_data) > 1:
            # Check which columns have ANY data across all rows
            cols_with_data = set()
            for row_idx, row in enumerate(table_data):
                for col_idx, cell in enumerate(row[:max_col]):
                    if cell and str(cell).strip():
                        cols_with_data.add(col_idx)
            
            # Keep all columns that have data somewhere
            cols_to_keep = sorted(list(cols_with_data))
            
            # Only filter if we're actually removing empty columns
            if cols_to_keep and len(cols_to_keep) < max_col:
                filtered_table_data = []
                for row in table_data:
                    filtered_row = [row[i] if i < len(row) else '' for i in cols_to_keep]
                    filtered_table_data.append(filtered_row)
                
                # Update merged cells info for new column indices
                filtered_merged_cells_info = []
                col_index_map = {old_idx: new_idx for new_idx, old_idx in enumerate(cols_to_keep)}
                
                for merge_info in merged_cells_info:
                    old_col = merge_info['col']
                    if old_col in col_index_map:
                        # Calculate new colspan based on how many columns in the span are kept
                        old_col_end = old_col + merge_info['colspan']
                        kept_cols_in_span = [c for c in range(old_col, old_col_end) if c in cols_to_keep]
                        
                        if kept_cols_in_span:
                            filtered_merged_cells_info.append({
                                'row': merge_info['row'],
                                'col': col_index_map[old_col],
                                'rowspan': merge_info['rowspan'],
                                'colspan': len(kept_cols_in_span)
                            })
                
                table_data = filtered_table_data
                merged_cells_info = filtered_merged_cells_info
                actual_col_count = len(cols_to_keep)
                logger.debug(f"Filtered table from {max_col} to {actual_col_count} columns (removed completely empty columns)")
            else:
                actual_col_count = max_col
        else:
            actual_col_count = max_col
        
        # Debug logging
        logger.debug(f"Extracted {len(table_data)} rows with {len(merged_cells_info)} merged cells")
        if len(table_data) > 0:
            logger.debug(f"First row: {table_data[0][:10] if len(table_data[0]) > 10 else table_data[0]}")

        # Create table info for complex table handling
        if merged_cells_info:
            table_info = {
                'is_complex': True,
                'merged_cells': merged_cells_info,
                'row_count': len(table_data),
                'col_count': actual_col_count,  # Use filtered column count
                'cell_spans': {}
            }

            # Add span info
            for merge_info in merged_cells_info:
                key = (merge_info['row'], merge_info['col'])
                table_info['cell_spans'][key] = (merge_info['rowspan'], merge_info['colspan'])

            # Use HTML converter for complex tables
            return self._convert_table_to_html(table_data, table_info)
        else:
            # Filter out completely empty rows before converting to markdown
            filtered_data = []
            for idx, row in enumerate(table_data):
                if idx == 0 or any(cell.strip() for cell in row if cell):  # Keep headers or non-empty rows
                    filtered_data.append(row)
            
            # Use standard markdown conversion
            return self._convert_table_to_markdown(filtered_data) if filtered_data else ""

    def _dataframe_to_markdown(self, df: pd.DataFrame) -> str:
        """Convert pandas DataFrame to markdown table"""
        # Convert DataFrame to list of lists
        table_data = [df.columns.tolist()]  # Header
        table_data.extend(df.values.tolist())  # Data

        # Handle NaN values
        for i, row in enumerate(table_data):
            table_data[i] = [str(cell) if pd.notna(cell) else "" for cell in row]

        return self._convert_table_to_markdown(table_data)

    def _extract_images_from_sheet(self, sheet, sheet_num: int, ocr_images: bool,
                                   ocr_results_map: Dict[str, str] = {},
                                   image_data_cache: Dict[tuple, bytes] = {}) -> List[Dict[str, Any]]:
        """Extract images from an Excel sheet"""
        images = []
        sheet_name = sheet.title
        
        # First try standard image extraction
        for img_idx, image in enumerate(sheet._images):
            try:
                # Try to get cached image data first to avoid double _data() call
                cache_key = (sheet_name, img_idx)  # Use sheet_name for consistency
                if cache_key in image_data_cache:
                    image_data = image_data_cache[cache_key]
                else:
                    # Fallback to extracting if not cached (shouldn't happen)
                    try:
                        image_data = image._data()
                        # Cache it now to avoid future calls
                        image_data_cache[cache_key] = image_data
                    except Exception as e:
                        logger.warning(f"Failed to extract image data: {e}")
                        continue

                if ocr_images:
                    # Use image content hash to find OCR result
                    img_hash = hash(image_data)

                    if img_hash in ocr_results_map:
                        ocr_text = ocr_results_map[img_hash]
                        images.append({
                            "type": "text:image_description",
                            "content": f"<image_ocr_result>{ocr_text}</image_ocr_result>",
                            "page": sheet_num
                        })
                    else:
                        # Fallback to individual OCR
                        logger.warning(f"OCR result not found for image on sheet {sheet_name}, using fallback OCR")
                        ocr_text = self._ocr_image(image_data)
                        images.append({
                            "type": "text:image_description",
                            "content": f"<image_ocr_result>{ocr_text}</image_ocr_result>",
                            "page": sheet_num
                        })
                else:
                    # Return base64 encoded image
                    base64_data = self._extract_image_as_base64(image_data)
                    images.append({
                        "type": "image",
                        "content": base64_data,
                        "page": sheet_num
                    })

            except Exception as e:
                logger.warning(f"Failed to extract image: {e}")
        
        # If no images found via standard method and this is the first sheet, check for fallback images
        if len(images) == 0 and sheet_num == 1:
            # Check if we have fallback images in cache
            fallback_idx = 0
            while True:
                fallback_key = ('_fallback', fallback_idx)
                if fallback_key in image_data_cache:
                    image_data = image_data_cache[fallback_key]
                    
                    if ocr_images:
                        # First try the special fallback OCR key for this specific image
                        fallback_ocr_key = ('_fallback_ocr', fallback_idx)
                        
                        if fallback_ocr_key in ocr_results_map:
                            ocr_text = ocr_results_map[fallback_ocr_key]
                            images.append({
                                "type": "text:image_description",
                                "content": f"<image_ocr_result>{ocr_text}</image_ocr_result>",
                                "page": sheet_num
                            })
                            logger.info(f"Using OCR result for fallback image {fallback_idx}")
                        else:
                            # Try hash-based lookup as backup
                            img_hash = hash(image_data)
                            if img_hash in ocr_results_map:
                                ocr_text = ocr_results_map[img_hash]
                                images.append({
                                    "type": "text:image_description",
                                    "content": f"<image_ocr_result>{ocr_text}</image_ocr_result>",
                                    "page": sheet_num
                                })
                            else:
                                # Fallback to individual OCR
                                logger.warning(f"OCR result not found for fallback image {fallback_idx}, using fallback OCR")
                                ocr_text = self._ocr_image(image_data)
                                images.append({
                                    "type": "text:image_description",
                                    "content": f"<image_ocr_result>{ocr_text}</image_ocr_result>",
                                    "page": sheet_num
                                })
                    else:
                        # Return base64 encoded image
                        base64_data = self._extract_image_as_base64(image_data)
                        images.append({
                            "type": "image",
                            "content": base64_data,
                            "page": sheet_num
                        })
                    
                    logger.info(f"Added fallback image {fallback_idx} to sheet {sheet_name}")
                    fallback_idx += 1
                else:
                    break  # No more fallback images

        return images

    def _collect_all_images(self) -> List[Dict[str, Any]]:
        """Collect all images from XLSX workbook for batch processing"""
        images_info = []
        image_counter = 0

        # First try standard openpyxl method
        for sheet_idx, sheet_name in enumerate(self.doc.sheetnames):
            sheet = self.doc[sheet_name]

            # Extract images from sheet using standard method
            for img_idx, image in enumerate(sheet._images):
                try:
                    image_data = image._data()
                    image_counter += 1
                    images_info.append({
                        'id': f'xlsx_sheet{sheet_idx + 1}_{sheet_name}_img{img_idx}',
                        'data': image_data,
                        'location': {'sheet': sheet_idx + 1, 'sheet_name': sheet_name, 'img_idx': img_idx}
                    })
                except Exception as e:
                    logger.warning(f"Failed to extract image from sheet {sheet_name}: {e}")

        # If no images found via standard method, try fallback ZIP extraction
        if len(images_info) == 0:
            logger.info("No images found via standard method, trying ZIP extraction fallback...")
            images_info = self._extract_images_from_zip()

        logger.info(f"Collected {len(images_info)} images from XLSX")
        return images_info
    
    def _extract_images_from_zip(self) -> List[Dict[str, Any]]:
        """Fallback method to extract images directly from XLSX ZIP structure
        
        This handles edge cases where images are embedded in non-standard ways:
        - Cell backgrounds
        - VML drawings
        - Legacy formats
        """
        images_info = []
        
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            
            with zipfile.ZipFile(self.file_path, 'r') as zip_file:
                # Look for all media files
                media_files = [f for f in zip_file.namelist() if '/media/' in f and 
                             any(f.lower().endswith(ext) for ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp'])]
                
                # Try to map images to their cell locations via drawing files
                image_to_cell_map = {}
                drawing_files = [f for f in zip_file.namelist() if '/drawings/' in f and f.endswith('.xml')]
                
                for drawing_file in drawing_files:
                    try:
                        content = zip_file.read(drawing_file).decode('utf-8')
                        # Parse drawing XML to find image anchors
                        # This is a simplified extraction - full implementation would need proper namespace handling
                        import re
                        # Look for cell references in anchor tags
                        # Pattern to find from cell references
                        from_cells = re.findall(r'<xdr:from>.*?<xdr:col>(\d+)</xdr:col>.*?<xdr:row>(\d+)</xdr:row>', 
                                               content, re.DOTALL)
                        # Map drawing index to cell location
                        for idx, (col, row) in enumerate(from_cells):
                            if idx < len(media_files):
                                # Convert to Excel cell reference (0-based to 1-based)
                                cell_ref = f"{chr(65 + int(col))}{int(row) + 1}"  # Simple conversion for single letters
                                image_to_cell_map[media_files[idx]] = cell_ref
                    except Exception as e:
                        logger.debug(f"Could not parse drawing file {drawing_file}: {e}")
                
                for idx, media_file in enumerate(media_files):
                    try:
                        image_data = zip_file.read(media_file)
                        
                        # Try to determine which sheet this image belongs to
                        # For now, we'll associate with the first sheet as a fallback
                        sheet_idx = 0
                        sheet_name = self.doc.sheetnames[0] if self.doc.sheetnames else "Sheet1"
                        
                        # Get cell reference if available
                        cell_ref = image_to_cell_map.get(media_file, None)
                        
                        images_info.append({
                            'id': f'xlsx_fallback_{idx}_{media_file.split("/")[-1]}',
                            'data': image_data,
                            'location': {
                                'sheet': sheet_idx + 1, 
                                'sheet_name': sheet_name, 
                                'img_idx': idx, 
                                'source': 'zip_fallback',
                                'cell_ref': cell_ref  # Add cell reference if found
                            },
                            'media_file': media_file  # Keep track of which media file this is
                        })
                        
                        logger.info(f"Extracted image via ZIP fallback: {media_file} (cell: {cell_ref or 'unknown'})")
                    except Exception as e:
                        logger.warning(f"Failed to extract {media_file}: {e}")
                        
        except Exception as e:
            logger.warning(f"ZIP fallback extraction failed: {e}")
            
        return images_info


class UniversalOfficeLoader:
    """Universal loader that detects file type and uses appropriate loader"""

    @staticmethod
    def load(file_path: Union[str, Path],
             extract_images: bool = True,
             ocr_images: bool = False,
             show_progress: bool = True,
             ocr=None,
             table_style: Union[str, TableStyle] = None) -> Dict[str, Any]:
        """Load any supported Office document
        
        Args:
            file_path: Path to the Office document
            extract_images: Whether to extract images
            ocr_images: Whether to OCR images
            show_progress: Whether to show progress logs
            ocr: OCR instance for image processing
            table_style: Output style for complex tables:
                - 'minimal_html': Clean HTML with only rowspan/colspan (default)
                - 'markdown_grid': Markdown with merge annotations
                - 'styled_html': Full HTML with inline styles (legacy)
        """
        file_path = Path(file_path)
        extension = file_path.suffix.lower()

        loaders = {
            '.docx': DocxLoader,
            '.pptx': PptxLoader,
            '.xlsx': XlsxLoader
        }

        if extension not in loaders:
            raise ValueError(f"Unsupported file type: {extension}")

        loader_class = loaders[extension]
        loader = loader_class(file_path, ocr=ocr, table_style=table_style)

        try:
            return loader.convert_to_json(
                extract_images=extract_images,
                ocr_images=ocr_images,
                show_progress=show_progress
            )
        finally:
            # Cleanup if needed
            pass


# Convenience functions
def office_to_json(
        file_path: Union[str, Path],
        output_path: Optional[Union[str, Path]] = None,
        output_markdown: bool = False,
        extract_images: bool = True,
        ocr_images: bool = False,
        show_progress: bool = True,
        ocr=None
) -> Dict[str, Any]:
    """
    Convert Office document to JSON with content in reading order
    
    Args:
        file_path: Path to the Office file (DOCX, PPTX, or XLSX)
        output_path: Optional path to save JSON output
        output_markdown: Also save as markdown file
        extract_images: Extract images as base64
        ocr_images: Use OCR to convert images to text descriptions
        show_progress: Show progress messages
    
    Returns:
        JSON data with content array
    """
    json_data = UniversalOfficeLoader.load(
        file_path,
        extract_images=extract_images,
        ocr_images=ocr_images,
        show_progress=show_progress,
        ocr=ocr
    )

    if output_path:
        output_path = Path(output_path)
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(json_data, f, ensure_ascii=False, indent=2)
        logger.info(f"JSON saved to: {output_path}")

        if output_markdown:
            markdown_path = output_path.with_suffix('.md')
            markdown_content = office_to_markdown(json_data)
            with open(markdown_path, 'w', encoding='utf-8') as f:
                f.write(markdown_content)
            logger.info(f"Markdown saved to: {markdown_path}")

    return json_data


def office_to_markdown(json_data: Dict[str, Any]) -> str:
    """Convert Office JSON data to markdown string"""
    markdown_parts = []

    for item in json_data["content"]:
        if item["type"] == "text:title":
            # Add titles with H1 formatting and extra spacing
            markdown_parts.append(f"# {item['content']}\n")
        elif item["type"] == "text:section":
            # Add sections with H2 formatting
            markdown_parts.append(f"## {item['content']}\n")
        elif item["type"] == "text:list":
            # Add list items with proper formatting
            markdown_parts.append(f"{item['content']}\n")
        elif item["type"] == "text:caption":
            # Add captions in italics
            markdown_parts.append(f"*{item['content']}*\n")
        elif item["type"] == "text:normal":
            # Add normal text with paragraph spacing
            markdown_parts.append(f"{item['content']}\n")
        elif item["type"] == "text:image_description":
            # Handle OCR results with XML tags
            ocr_text = item['content']
            if ocr_text.startswith('<image_ocr_result>') and ocr_text.endswith('</image_ocr_result>'):
                ocr_text = ocr_text[18:-19]  # Remove tags
            
            markdown_parts.append("```")
            markdown_parts.append("<ocr_result>")
            markdown_parts.append(ocr_text)
            markdown_parts.append("</ocr_result>")
            markdown_parts.append("```\n")
        elif item["type"] == "table":
            # Tables already have proper spacing
            markdown_parts.append(item["content"])
        elif item["type"] == "image":
            # Include image as markdown with base64 data URL
            markdown_parts.append(f'![Image](data:image/png;base64,{item["content"]})\n')

    return "\n".join(markdown_parts)
