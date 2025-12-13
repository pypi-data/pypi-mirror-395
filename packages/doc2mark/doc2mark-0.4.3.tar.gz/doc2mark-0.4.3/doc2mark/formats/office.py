"""Office format processors (DOCX, XLSX, PPTX)."""

import logging
import zipfile
import base64
from pathlib import Path
from typing import Any, Dict, List, Optional, Union, Tuple

from doc2mark.core.base import (
    BaseProcessor,
    DocumentFormat,
    DocumentMetadata,
    ProcessedDocument,
    ProcessingError
)
from doc2mark.ocr.base import BaseOCR

# Import the advanced pipeline loader
try:
    from doc2mark.pipelines.office_advanced_pipeline import (
        DocxLoader, PptxLoader, XlsxLoader, UniversalOfficeLoader
    )
    ADVANCED_PIPELINE_AVAILABLE = True
except ImportError:
    ADVANCED_PIPELINE_AVAILABLE = False
    logging.warning("Advanced Office pipeline not available. Using basic processing.")

logger = logging.getLogger(__name__)


class OfficeProcessor(BaseProcessor):
    """Processor for modern Office formats (DOCX, XLSX, PPTX) with advanced features."""

    def __init__(self, ocr: Optional[BaseOCR] = None, table_style: Optional[str] = None):
        """Initialize Office processor.
        
        Args:
            ocr: OCR provider for image extraction
            table_style: Output style for complex tables:
                - 'minimal_html': Clean HTML with only rowspan/colspan (default)
                - 'markdown_grid': Markdown with merge annotations  
                - 'styled_html': Full HTML with inline styles (legacy)
        """
        self.ocr = ocr
        self.table_style = table_style
        self._docx = None
        self._openpyxl = None
        self._pptx = None

    @property
    def python_docx(self):
        """Lazy load python-docx."""
        if self._docx is None:
            try:
                import docx
                self._docx = docx
            except ImportError:
                raise ImportError(
                    "python-docx is not installed. "
                    "Install it with: pip install python-docx"
                )
        return self._docx

    @property
    def openpyxl(self):
        """Lazy load openpyxl."""
        if self._openpyxl is None:
            try:
                import openpyxl
                self._openpyxl = openpyxl
            except ImportError:
                raise ImportError(
                    "openpyxl is not installed. "
                    "Install it with: pip install openpyxl"
                )
        return self._openpyxl

    @property
    def python_pptx(self):
        """Lazy load python-pptx."""
        if self._pptx is None:
            try:
                import pptx
                self._pptx = pptx
            except ImportError:
                raise ImportError(
                    "python-pptx is not installed. "
                    "Install it with: pip install python-pptx"
                )
        return self._pptx

    def can_process(self, file_path: Union[str, Path]) -> bool:
        """Check if this processor can handle the file."""
        file_path = Path(file_path)
        extension = file_path.suffix.lower().lstrip('.')
        return extension in ['docx', 'xlsx', 'pptx']

    def process(
            self,
            file_path: Union[str, Path],
            **kwargs
    ) -> ProcessedDocument:
        """Process Office document using advanced pipeline if available."""
        file_path = Path(file_path)
        extension = file_path.suffix.lower().lstrip('.')

        # Get file size
        file_size = file_path.stat().st_size

        # Use advanced pipeline if available
        if ADVANCED_PIPELINE_AVAILABLE:
            content, metadata, images = self._process_with_advanced_pipeline(file_path, **kwargs)
        else:
            # Fallback to basic processing
            if extension == 'docx':
                content, metadata, images = self._process_docx_basic(file_path, **kwargs)
                doc_format = DocumentFormat.DOCX
            elif extension == 'xlsx':
                content, metadata, images = self._process_xlsx_basic(file_path, **kwargs)
                doc_format = DocumentFormat.XLSX
            elif extension == 'pptx':
                content, metadata, images = self._process_pptx_basic(file_path, **kwargs)
                doc_format = DocumentFormat.PPTX
            else:
                raise ProcessingError(f"Unsupported Office format: {extension}")

        # Determine format
        if extension == 'docx':
            doc_format = DocumentFormat.DOCX
        elif extension == 'xlsx':
            doc_format = DocumentFormat.XLSX
        elif extension == 'pptx':
            doc_format = DocumentFormat.PPTX

        # Build metadata
        doc_metadata = DocumentMetadata(
            filename=file_path.name,
            format=doc_format,
            size_bytes=file_size,
            **metadata
        )

        return ProcessedDocument(
            content=content,
            metadata=doc_metadata,
            images=images
        )

    def _process_with_advanced_pipeline(
        self, 
        file_path: Path, 
        **kwargs
    ) -> Tuple[str, dict, List[Dict[str, Any]]]:
        """Process document using advanced pipeline with all features."""
        try:
            # Configure options
            extract_images = kwargs.get('extract_images', False)
            ocr_images = extract_images and self.ocr is not None
            
            # Use the advanced pipeline
            json_data = UniversalOfficeLoader.load(
                file_path,
                extract_images=extract_images,
                ocr_images=ocr_images,
                show_progress=kwargs.get('show_progress', False),
                ocr=self.ocr,
                table_style=kwargs.get('table_style', self.table_style)
            )
            
            # Convert JSON data to markdown
            markdown_parts = []
            images = []
            
            for item in json_data["content"]:
                if item["type"] == "text:title":
                    markdown_parts.append(f"# {item['content']}\n")
                elif item["type"] == "text:section":
                    markdown_parts.append(f"## {item['content']}\n")
                elif item["type"] == "text:list":
                    markdown_parts.append(f"{item['content']}\n")
                elif item["type"] == "text:caption":
                    markdown_parts.append(f"*{item['content']}*\n")
                elif item["type"] == "text:normal":
                    markdown_parts.append(f"{item['content']}\n")
                elif item["type"] == "text:image_description":
                    # Handle OCR results
                    ocr_text = item['content']
                    if ocr_text.startswith('<image_ocr_result>') and ocr_text.endswith('</image_ocr_result>'):
                        ocr_text = ocr_text[18:-19]  # Remove tags
                    
                    # Format OCR result in XML tags within markdown code block
                    markdown_parts.append("```xml")
                    markdown_parts.append("<ocr_result>")
                    markdown_parts.append(ocr_text)
                    markdown_parts.append("</ocr_result>")
                    markdown_parts.append("```\n")
                elif item["type"] == "table":
                    # Tables already include proper formatting (markdown or HTML)
                    markdown_parts.append(item["content"])
                elif item["type"] == "image":
                    # Store image data
                    images.append({
                        'data': base64.b64decode(item["content"]),
                        'page': item.get('page', 1)
                    })
                    markdown_parts.append(f"![Image {len(images)}]\n")
            
            content = '\n'.join(markdown_parts)
            
            # Build metadata
            metadata = {
                'page_count': json_data.get('pages', 1),
                'image_count': len(images),
            }
            
            # Add format-specific metadata
            if file_path.suffix.lower() == '.docx':
                metadata['word_count'] = len(content.split())
            elif file_path.suffix.lower() == '.xlsx':
                # XLSX specific metadata is already included in json_data
                pass
            elif file_path.suffix.lower() == '.pptx':
                # Count slides from content
                slide_count = content.count('Slide ')
                metadata['slide_count'] = max(slide_count, 1)
            
            return content, metadata, images
            
        except Exception as e:
            logger.warning(f"Advanced pipeline processing failed: {e}, falling back to basic processing")
            # Fallback to basic processing
            extension = file_path.suffix.lower().lstrip('.')
            if extension == 'docx':
                return self._process_docx_basic(file_path, **kwargs)
            elif extension == 'xlsx':
                return self._process_xlsx_basic(file_path, **kwargs)
            elif extension == 'pptx':
                return self._process_pptx_basic(file_path, **kwargs)
            else:
                raise ProcessingError(f"Unsupported Office format: {extension}")

    def _process_docx_basic(self, file_path: Path, **kwargs) -> Tuple[str, dict, List[Dict[str, Any]]]:
        """Basic DOCX processing (fallback when advanced pipeline not available)."""
        try:
            doc = self.python_docx.Document(str(file_path))

            # Extract text content
            markdown_parts = []

            # Process paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    # Check style for headings
                    if para.style.name.startswith('Heading'):
                        level = int(para.style.name[-1]) if para.style.name[-1].isdigit() else 1
                        markdown_parts.append(f"{'#' * level} {para.text}")
                    else:
                        # Handle text formatting
                        text = self._format_docx_paragraph(para)
                        markdown_parts.append(text)
                    markdown_parts.append("")

            # Process tables
            for table in doc.tables:
                table_md = self._convert_docx_table_to_markdown_basic(table)
                markdown_parts.append(table_md)
                markdown_parts.append("")

            # Extract images if requested
            images = []
            if kwargs.get('extract_images', False) and self.ocr:
                images = self._extract_docx_images_basic(file_path)

            # Metadata
            metadata = {
                'page_count': len(doc.element.xpath('//w:sectPr')),
                'word_count': sum(len(para.text.split()) for para in doc.paragraphs),
                'author': doc.core_properties.author,
                'title': doc.core_properties.title,
                'creation_date': str(doc.core_properties.created) if doc.core_properties.created else None,
                'modification_date': str(doc.core_properties.modified) if doc.core_properties.modified else None,
                'image_count': len(images),
            }

            return '\n'.join(markdown_parts), metadata, images

        except Exception as e:
            logger.error(f"Failed to process DOCX: {e}")
            raise ProcessingError(f"DOCX processing failed: {str(e)}")

    def _process_xlsx_basic(self, file_path: Path, **kwargs) -> Tuple[str, dict, List[Dict[str, Any]]]:
        """Basic XLSX processing (fallback when advanced pipeline not available)."""
        try:
            wb = self.openpyxl.load_workbook(str(file_path), data_only=True)

            markdown_parts = []
            total_cells = 0

            # Process each sheet
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                
                # Find actual max column with data (not Excel's 16384 limit)
                actual_max_col = 0
                actual_max_row = 0
                
                # Check first 100 rows to find actual data extent
                for row_idx, row in enumerate(sheet.iter_rows(min_row=1, max_row=min(sheet.max_row, 100)), 1):
                    for col_idx, cell in enumerate(row, 1):
                        if cell.value is not None:
                            actual_max_col = max(actual_max_col, col_idx)
                            actual_max_row = max(actual_max_row, row_idx)
                
                # Check rest of rows if needed
                if sheet.max_row > 100:
                    for row in sheet.iter_rows(min_row=101, max_row=sheet.max_row, max_col=actual_max_col):
                        if any(cell.value is not None for cell in row):
                            actual_max_row = sheet.max_row
                            break

                if actual_max_row > 0 and actual_max_col > 0:
                    markdown_parts.append(f"## {sheet_name}")
                    markdown_parts.append("")

                    # Convert sheet to markdown table
                    table_data = []
                    for row in sheet.iter_rows(min_row=1, max_row=actual_max_row, 
                                               min_col=1, max_col=actual_max_col, 
                                               values_only=True):
                        # Filter out completely empty rows
                        if any(cell is not None for cell in row):
                            table_data.append([str(cell) if cell is not None else "" for cell in row])
                            total_cells += actual_max_col

                    if table_data:
                        table_md = self._convert_list_to_markdown_table(table_data)
                        markdown_parts.append(table_md)
                        markdown_parts.append("")

            # Extract images if requested
            images = []
            if kwargs.get('extract_images', False) and self.ocr:
                images = self._extract_xlsx_images_basic(file_path)
                
                # Add extracted image text to content if available
                if images:
                    image_texts = [img['text'] for img in images if img.get('text', '').strip()]
                    if image_texts:
                        markdown_parts.append("## Extracted Images")
                        markdown_parts.append("")
                        for i, text in enumerate(image_texts, 1):
                            markdown_parts.append(f"### Image {i}")
                            markdown_parts.append("```xml")
                            markdown_parts.append("<ocr_result>")
                            markdown_parts.append(text)
                            markdown_parts.append("</ocr_result>")
                            markdown_parts.append("```")
                            markdown_parts.append("")

            # Metadata
            metadata = {
                'page_count': len(wb.sheetnames),  # Using page_count to represent number of sheets
                'sheet_names': wb.sheetnames,
                'total_cells': total_cells,
                'image_count': len(images),
            }

            return '\n'.join(markdown_parts), metadata, images

        except Exception as e:
            logger.error(f"Failed to process XLSX: {e}")
            raise ProcessingError(f"XLSX processing failed: {str(e)}")

    def _process_pptx_basic(self, file_path: Path, **kwargs) -> Tuple[str, dict, List[Dict[str, Any]]]:
        """Basic PPTX processing (fallback when advanced pipeline not available)."""
        try:
            prs = self.python_pptx.Presentation(str(file_path))

            markdown_parts = []
            image_count = 0

            # Process each slide
            for i, slide in enumerate(prs.slides, 1):
                markdown_parts.append(f"## Slide {i}")
                markdown_parts.append("")

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Handle title shapes
                        if shape == slide.shapes.title:
                            markdown_parts.append(f"### {shape.text}")
                        else:
                            markdown_parts.append(shape.text)
                        markdown_parts.append("")

                    # Count images
                    if shape.shape_type == 13:  # Picture
                        image_count += 1

                # Process tables
                for shape in slide.shapes:
                    if shape.has_table:
                        table_md = self._convert_pptx_table_to_markdown_basic(shape.table)
                        markdown_parts.append(table_md)
                        markdown_parts.append("")

            # Extract images if requested
            images = []
            if kwargs.get('extract_images', False) and self.ocr:
                images = self._extract_pptx_images_basic(file_path)
                
                # Add extracted image text to content if available
                if images:
                    image_texts = [img['text'] for img in images if img.get('text', '').strip()]
                    if image_texts:
                        markdown_parts.append("## Extracted Images")
                        markdown_parts.append("")
                        for i, text in enumerate(image_texts, 1):
                            markdown_parts.append(f"### Image {i}")
                            markdown_parts.append("```xml")
                            markdown_parts.append("<ocr_result>")
                            markdown_parts.append(text)
                            markdown_parts.append("</ocr_result>")
                            markdown_parts.append("```")
                            markdown_parts.append("")

            # Metadata
            metadata = {
                'page_count': len(prs.slides),
                'slide_count': len(prs.slides),
                'image_count': image_count,
                'title': prs.core_properties.title,
                'author': prs.core_properties.author,
            }

            return '\n'.join(markdown_parts), metadata, images

        except Exception as e:
            logger.error(f"Failed to process PPTX: {e}")
            raise ProcessingError(f"PPTX processing failed: {str(e)}")

    def _format_docx_paragraph(self, paragraph) -> str:
        """Format DOCX paragraph with inline formatting."""
        if not paragraph.runs:
            return paragraph.text

        formatted_text = []
        for run in paragraph.runs:
            text = run.text
            if run.bold:
                text = f"**{text}**"
            if run.italic:
                text = f"*{text}*"
            formatted_text.append(text)

        return ''.join(formatted_text)

    def _convert_docx_table_to_markdown_basic(self, table) -> str:
        """Convert DOCX table to markdown (basic version)."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)

        return self._convert_list_to_markdown_table(rows)

    def _convert_pptx_table_to_markdown_basic(self, table) -> str:
        """Convert PPTX table to markdown (basic version)."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)

        return self._convert_list_to_markdown_table(rows)

    def _convert_list_to_markdown_table(self, data: List[List[str]]) -> str:
        """Convert list of lists to markdown table."""
        if not data:
            return ""

        # Determine column widths
        col_widths = [0] * len(data[0])
        for row in data:
            for i, cell in enumerate(row):
                col_widths[i] = max(col_widths[i], len(cell))

        # Build table
        lines = []

        # Header
        header = "| " + " | ".join(cell.ljust(col_widths[i]) for i, cell in enumerate(data[0])) + " |"
        lines.append(header)

        # Separator
        separator = "|" + "|".join("-" * (w + 2) for w in col_widths) + "|"
        lines.append(separator)

        # Data rows
        for row in data[1:]:
            row_str = "| " + " | ".join(cell.ljust(col_widths[i]) for i, cell in enumerate(row)) + " |"
            lines.append(row_str)

        return '\n'.join(lines)

    def _extract_docx_images_basic(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract images from DOCX file using basic method."""
        images = []

        try:
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                # Collect all images for batch processing
                image_batch = []
                image_info = []

                # Find image files in the media folder
                for file_info in zip_file.filelist:
                    if file_info.filename.startswith('word/media/'):
                        image_data = zip_file.read(file_info.filename)
                        image_batch.append(image_data)
                        image_info.append({
                            'filename': file_info.filename,
                            'size': len(image_data)
                        })

                # Batch process OCR if we have images and OCR is available
                if image_batch and self.ocr:
                    try:
                        # Prepare OCR kwargs with language configuration if available
                        ocr_kwargs = {}
                        if hasattr(self.ocr, 'config') and self.ocr.config and self.ocr.config.language:
                            ocr_kwargs['language'] = self.ocr.config.language
                            logger.info(
                                f"🌍 Passing language configuration to Office batch OCR: {self.ocr.config.language}")

                        logger.info(f"🚀 Processing {len(image_batch)} DOCX images with batch OCR")
                        ocr_results = self.ocr.batch_process_images(image_batch, **ocr_kwargs)

                        # Combine results
                        for info, ocr_result in zip(image_info, ocr_results):
                            text = ocr_result.text if hasattr(ocr_result, 'text') else str(ocr_result)
                            info['text'] = text
                            images.append(info)

                        logger.info(f"✅ DOCX batch OCR completed successfully")

                    except Exception as e:
                        logger.warning(f"Failed to batch OCR DOCX images: {e}")
                        # Add images without OCR text
                        for info in image_info:
                            info['text'] = ''
                            images.append(info)
                else:
                    # No OCR available, add images without text
                    for info in image_info:
                        info['text'] = ''
                        images.append(info)

        except Exception as e:
            logger.warning(f"Failed to extract images from DOCX: {e}")

        return images

    def _extract_xlsx_images_basic(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract images from XLSX file using basic method."""
        images = []

        try:
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                # Collect all images for batch processing
                image_batch = []
                image_info = []

                # Find image files in the media folder
                for file_info in zip_file.filelist:
                    if file_info.filename.startswith('xl/media/'):
                        image_data = zip_file.read(file_info.filename)
                        image_batch.append(image_data)
                        image_info.append({
                            'filename': file_info.filename,
                            'size': len(image_data)
                        })

                # Batch process OCR if we have images and OCR is available
                if image_batch and self.ocr:
                    try:
                        # Prepare OCR kwargs with language configuration if available
                        ocr_kwargs = {}
                        if hasattr(self.ocr, 'config') and self.ocr.config and self.ocr.config.language:
                            ocr_kwargs['language'] = self.ocr.config.language
                            logger.info(
                                f"🌍 Passing language configuration to Office batch OCR: {self.ocr.config.language}")

                        logger.info(f"🚀 Processing {len(image_batch)} XLSX images with batch OCR")
                        ocr_results = self.ocr.batch_process_images(image_batch, **ocr_kwargs)

                        # Combine results
                        for info, ocr_result in zip(image_info, ocr_results):
                            text = ocr_result.text if hasattr(ocr_result, 'text') else str(ocr_result)
                            info['text'] = text
                            images.append(info)

                        logger.info(f"✅ XLSX batch OCR completed successfully")

                    except Exception as e:
                        logger.warning(f"Failed to batch OCR XLSX images: {e}")
                        # Add images without OCR text
                        for info in image_info:
                            info['text'] = ''
                            images.append(info)
                else:
                    # No OCR available, add images without text
                    for info in image_info:
                        info['text'] = ''
                        images.append(info)

        except Exception as e:
            logger.warning(f"Failed to extract images from XLSX: {e}")

        return images

    def _extract_pptx_images_basic(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract images from PPTX file using basic method."""
        images = []

        try:
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                # Collect all images for batch processing
                image_batch = []
                image_info = []

                # Find image files in the media folder
                for file_info in zip_file.filelist:
                    if file_info.filename.startswith('ppt/media/'):
                        image_data = zip_file.read(file_info.filename)
                        image_batch.append(image_data)
                        image_info.append({
                            'filename': file_info.filename,
                            'size': len(image_data)
                        })

                # Batch process OCR if we have images and OCR is available
                if image_batch and self.ocr:
                    try:
                        # Prepare OCR kwargs with language configuration if available
                        ocr_kwargs = {}
                        if hasattr(self.ocr, 'config') and self.ocr.config and self.ocr.config.language:
                            ocr_kwargs['language'] = self.ocr.config.language
                            logger.info(
                                f"🌍 Passing language configuration to Office batch OCR: {self.ocr.config.language}")

                        logger.info(f"🚀 Processing {len(image_batch)} PPTX images with batch OCR")
                        ocr_results = self.ocr.batch_process_images(image_batch, **ocr_kwargs)

                        # Combine results
                        for info, ocr_result in zip(image_info, ocr_results):
                            text = ocr_result.text if hasattr(ocr_result, 'text') else str(ocr_result)
                            info['text'] = text
                            images.append(info)

                        logger.info(f"✅ PPTX batch OCR completed successfully")

                    except Exception as e:
                        logger.warning(f"Failed to batch OCR PPTX images: {e}")
                        # Add images without OCR text
                        for info in image_info:
                            info['text'] = ''
                            images.append(info)
                else:
                    # No OCR available, add images without text
                    for info in image_info:
                        info['text'] = ''
                        images.append(info)

        except Exception as e:
            logger.warning(f"Failed to extract images from PPTX: {e}")

        return images
