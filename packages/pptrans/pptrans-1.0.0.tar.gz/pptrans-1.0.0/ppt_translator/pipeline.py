"""PowerPoint translation pipeline utilities."""
from __future__ import annotations

import json
import xml.etree.ElementTree as ET
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Optional, List
from xml.dom import minidom

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .translation import TranslationService
from .vision import VisionReviewer
from .render import render_slide_to_image
from .review import ReviewFileGenerator, ReviewFileLoader, SlideTranslation


def get_alignment_value(alignment_str: str | None):
    """Convert alignment string to PP_ALIGN enum value."""
    alignment_map = {
        "PP_ALIGN.CENTER": PP_ALIGN.CENTER,
        "PP_ALIGN.LEFT": PP_ALIGN.LEFT,
        "PP_ALIGN.RIGHT": PP_ALIGN.RIGHT,
        "PP_ALIGN.JUSTIFY": PP_ALIGN.JUSTIFY,
        "None": None,
        None: None,
    }
    return alignment_map.get(alignment_str)


def get_shape_properties(shape):
    """Extract text shape properties."""
    shape_data = {
        "text": "",
        "font_size": None,
        "font_name": None,
        "alignment": None,
        "width": shape.width,
        "height": shape.height,
        "left": shape.left,
        "top": shape.top,
        "bold": None,
        "italic": None,
        "line_spacing": None,
        "space_before": None,
        "space_after": None,
        "font_color": None,
    }
    if hasattr(shape, "text"):
        shape_data["text"] = shape.text.strip()
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.runs:
                    run = paragraph.runs[0]
                    if getattr(run.font, "size", None) is not None:
                        shape_data["font_size"] = run.font.size.pt
                    if getattr(run.font, "name", None):
                        shape_data["font_name"] = run.font.name
                    if hasattr(run.font, "bold"):
                        shape_data["bold"] = run.font.bold
                    if hasattr(run.font, "italic"):
                        shape_data["italic"] = run.font.italic
                    if (
                        getattr(run.font, "color", None) is not None
                        and getattr(run.font.color, "rgb", None) is not None
                    ):
                        shape_data["font_color"] = str(run.font.color.rgb)
                if getattr(paragraph, "line_spacing", None) is not None:
                    shape_data["line_spacing"] = paragraph.line_spacing
                if getattr(paragraph, "space_before", None) is not None:
                    shape_data["space_before"] = paragraph.space_before
                if getattr(paragraph, "space_after", None) is not None:
                    shape_data["space_after"] = paragraph.space_after
                if getattr(paragraph, "alignment", None) is not None:
                    shape_data["alignment"] = f"PP_ALIGN.{paragraph.alignment}" if paragraph.alignment else None
    return shape_data


def estimate_text_dimensions(text: str, font_size_pt: float, font_name: str = "Arial", width_emu: int = None) -> tuple[float, float]:
    """Estimate text dimensions in EMU (English Metric Units).
    
    Returns (estimated_width, estimated_height) in EMU.
    This is a rough estimation - actual rendering may vary.
    """
    # Approximate character width: font_size * 0.6 (for most fonts)
    # Approximate line height: font_size * 1.2
    char_width_emu = font_size_pt * 0.6 * 12700  # Convert pt to EMU (1pt = 12700 EMU)
    line_height_emu = font_size_pt * 1.2 * 12700
    
    if width_emu:
        # Calculate how many characters fit per line
        chars_per_line = max(1, int(width_emu / char_width_emu))
        # Calculate number of lines needed
        num_lines = max(1, (len(text) + chars_per_line - 1) // chars_per_line)
        estimated_width = min(width_emu, len(text) * char_width_emu)
        estimated_height = num_lines * line_height_emu
    else:
        # No width constraint, estimate single line
        estimated_width = len(text) * char_width_emu
        estimated_height = line_height_emu
    
    return (estimated_width, estimated_height)


def check_text_fits(text: str, font_size_pt: float, box_width_emu: int, box_height_emu: int, font_name: str = "Arial") -> tuple[bool, float]:
    """Check if text fits in the given box dimensions.
    
    Returns (fits, suggested_font_size).
    If text doesn't fit, suggests a smaller font size.
    """
    estimated_width, estimated_height = estimate_text_dimensions(text, font_size_pt, font_name, box_width_emu)
    
    fits = estimated_width <= box_width_emu and estimated_height <= box_height_emu
    
    if not fits:
        # Calculate scale factor needed
        width_scale = box_width_emu / estimated_width if estimated_width > 0 else 1.0
        height_scale = box_height_emu / estimated_height if estimated_height > 0 else 1.0
        scale_factor = min(width_scale, height_scale, 1.0)  # Don't scale up
        
        # Suggest new font size (with some margin)
        suggested_size = font_size_pt * scale_factor * 0.95  # 5% margin
        suggested_size = max(8.0, suggested_size)  # Minimum 8pt
    else:
        suggested_size = font_size_pt
    
    return (fits, suggested_size)


def apply_shape_properties(shape, shape_data, auto_adjust_font: bool = True):
    """Apply saved properties to a shape with optional layout-aware font adjustment."""
    try:
        shape.width = shape_data["width"]
        shape.height = shape_data["height"]
        shape.left = shape_data["left"]
        shape.top = shape_data["top"]
        shape.text = ""
        paragraph = shape.text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = shape_data["text"]
        
        original_font_size = shape_data.get("font_size") or 12.0
        font_size = original_font_size
        
        # Layout-aware adjustment: check if text fits
        if auto_adjust_font and shape_data["text"]:
            fits, suggested_size = check_text_fits(
                shape_data["text"],
                original_font_size,
                shape_data["width"],
                shape_data["height"],
                shape_data.get("font_name") or "Arial"
            )
            if not fits:
                font_size = suggested_size
                print(f"  ⚠️  Text overflow detected, adjusting font size: {original_font_size:.1f}pt → {font_size:.1f}pt")
        
        # Apply font size (with default scaling for translation)
        if font_size:
            adjusted_size = font_size * 0.7  # Default scaling
            run.font.size = Pt(adjusted_size)
        
        run.font.name = shape_data.get("font_name") or "Arial"
        if shape_data.get("font_color"):
            run.font.color.rgb = RGBColor.from_string(shape_data["font_color"])
        if shape_data.get("bold") is not None:
            run.font.bold = shape_data["bold"]
        if shape_data.get("italic") is not None:
            run.font.italic = shape_data["italic"]
        if shape_data.get("alignment"):
            paragraph.alignment = get_alignment_value(shape_data["alignment"])
        if shape_data.get("line_spacing"):
            paragraph.line_spacing = shape_data["line_spacing"]
        if shape_data.get("space_before"):
            paragraph.space_before = shape_data["space_before"]
        if shape_data.get("space_after"):
            paragraph.space_after = shape_data["space_after"]
    except Exception as exc:  # pragma: no cover - best effort logging
        print(f"Error applying shape properties: {exc}")


def get_table_properties(table):
    """Extract table properties."""
    table_data = {
        "rows": len(table.rows),
        "cols": len(table.columns),
        "cells": [],
    }
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            cell_data = {
                "text": cell.text.strip(),
                "font_size": None,
                "font_name": None,
                "alignment": None,
                "margin_left": cell.margin_left,
                "margin_right": cell.margin_right,
                "margin_top": cell.margin_top,
                "margin_bottom": cell.margin_bottom,
                "vertical_anchor": str(cell.vertical_anchor) if cell.vertical_anchor else None,
                "font_color": None,
            }
            if cell.text_frame.paragraphs:
                paragraph = cell.text_frame.paragraphs[0]
                if paragraph.runs:
                    run = paragraph.runs[0]
                    if getattr(run.font, "size", None) is not None:
                        cell_data["font_size"] = run.font.size.pt
                    if getattr(run.font, "name", None):
                        cell_data["font_name"] = run.font.name
                    if hasattr(run.font, "bold"):
                        cell_data["bold"] = run.font.bold
                    if hasattr(run.font, "italic"):
                        cell_data["italic"] = run.font.italic
                    if (
                        getattr(run.font, "color", None) is not None
                        and getattr(run.font.color, "rgb", None) is not None
                    ):
                        cell_data["font_color"] = str(run.font.color.rgb)
                if getattr(paragraph, "alignment", None) is not None:
                    cell_data["alignment"] = f"PP_ALIGN.{paragraph.alignment}" if paragraph.alignment else None
            row_data.append(cell_data)
        table_data["cells"].append(row_data)
    return table_data


def apply_table_properties(table, table_data):
    """Apply saved table properties."""
    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            try:
                cell_data = table_data["cells"][row_idx][col_idx]
                cell.margin_left = cell_data["margin_left"]
                cell.margin_right = cell_data["margin_right"]
                cell.margin_top = cell_data["margin_top"]
                cell.margin_bottom = cell_data["margin_bottom"]
                if cell_data.get("vertical_anchor"):
                    cell.vertical_anchor = eval(cell_data["vertical_anchor"])
                cell.text = ""
                paragraph = cell.text_frame.paragraphs[0]
                run = paragraph.add_run()
                run.text = cell_data["text"]
                if cell_data.get("font_size"):
                    adjusted_size = cell_data["font_size"] * 0.8
                    run.font.size = Pt(adjusted_size)
                run.font.name = cell_data.get("font_name") or "Arial"
                if cell_data.get("font_color"):
                    run.font.color.rgb = RGBColor.from_string(cell_data["font_color"])
                if "bold" in cell_data:
                    run.font.bold = cell_data["bold"]
                if "italic" in cell_data:
                    run.font.italic = cell_data["italic"]
                if cell_data.get("alignment"):
                    paragraph.alignment = get_alignment_value(cell_data["alignment"])
            except Exception as exc:  # pragma: no cover - best effort logging
                print(f"Error setting cell properties: {exc}")


def extract_text_from_slide(
    slide,
    slide_number: int,
    *,
    translator: TranslationService | None,
    source_lang: str,
    target_lang: str,
):
    """Extract text from a slide and optionally translate it."""
    slide_element = ET.Element("slide")
    slide_element.set("number", str(slide_number))
    for shape_index, shape in enumerate(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_element = ET.SubElement(slide_element, "table_element")
            table_element.set("shape_index", str(shape_index))
            table_data = get_table_properties(shape.table)
            if translator:
                for row in table_data["cells"]:
                    for cell in row:
                        cell["text"] = translator.translate(cell["text"], source_lang, target_lang)
            props_element = ET.SubElement(table_element, "properties")
            props_element.text = json.dumps(table_data, indent=2)
        elif hasattr(shape, "text"):
            text_element = ET.SubElement(slide_element, "text_element")
            text_element.set("shape_index", str(shape_index))
            shape_data = get_shape_properties(shape)
            if translator:
                shape_data["text"] = translator.translate(shape_data["text"], source_lang, target_lang)
            props_element = ET.SubElement(text_element, "properties")
            props_element.text = json.dumps(shape_data, indent=2)
    return slide_element


def ppt_to_xml(
    ppt_path: str,
    *,
    translator: TranslationService | None,
    source_lang: str,
    target_lang: str,
    max_workers: int = 4,
    temp_dir: Path | None = None,
) -> Optional[str]:
    """Convert a PowerPoint presentation to XML."""
    root = ET.Element("presentation")
    base_dir = Path(ppt_path).parent
    if temp_dir is None:
        temp_dir = base_dir
    try:
        prs = Presentation(ppt_path)
        root.set("file_path", Path(ppt_path).name)
        workers = max(1, max_workers)
        with ThreadPoolExecutor(max_workers=workers) as executor:
            future_to_slide = {
                executor.submit(
                    extract_text_from_slide,
                    slide,
                    slide_number,
                    translator=translator,
                    source_lang=source_lang,
                    target_lang=target_lang,
                ): slide_number
                for slide_number, slide in enumerate(prs.slides, start=1)
            }
            for future, slide_number in future_to_slide.items():
                slide_element = future.result()
                root.append(slide_element)
                intermediate_path = temp_dir / f"slide_{slide_number}_{'translated' if translator else 'original'}.xml"
                xml_str = minidom.parseString(ET.tostring(root)).toprettyxml(indent="  ")
                with open(intermediate_path, "w", encoding="utf-8") as handle:
                    handle.write(xml_str)
        return minidom.parseString(ET.tostring(root)).toprettyxml(indent="  ")
    except Exception as exc:  # pragma: no cover - best effort logging
        print(f"Error processing presentation: {exc}")
        return None


def create_translated_ppt(original_ppt_path: str, translated_xml_path: str, output_ppt_path: str) -> None:
    """Create a new PowerPoint presentation using translated content."""
    try:
        prs = Presentation(original_ppt_path)
        tree = ET.parse(translated_xml_path)
        root = tree.getroot()
        for slide_number, slide in enumerate(prs.slides, start=1):
            xml_slide = root.find(f".//slide[@number='{slide_number}']")
            if xml_slide is None:
                continue
            for shape_index, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_element = xml_slide.find(f".//table_element[@shape_index='{shape_index}']")
                    if table_element is not None:
                        props_element = table_element.find("properties")
                        if props_element is not None and props_element.text:
                            try:
                                table_data = json.loads(props_element.text)
                                apply_table_properties(shape.table, table_data)
                            except Exception as exc:  # pragma: no cover
                                print(f"Error applying table properties: {exc}")
                elif hasattr(shape, "text"):
                    text_element = xml_slide.find(f".//text_element[@shape_index='{shape_index}']")
                    if text_element is not None:
                        props_element = text_element.find("properties")
                        if props_element is not None and props_element.text:
                            try:
                                shape_data = json.loads(props_element.text)
                                apply_shape_properties(shape, shape_data, auto_adjust_font=True)
                            except Exception as exc:  # pragma: no cover
                                print(f"Error applying shape properties: {exc}")
        prs.save(output_ppt_path)
        print(f"Translated PowerPoint saved to: {output_ppt_path}")
    except Exception as exc:  # pragma: no cover - logging only
        print(f"Error creating translated PowerPoint: {exc}")


def _collect_review_data(
    original_xml_path: Optional[Path],
    translated_xml_path: Path,
    vision_reviewer: Optional[VisionReviewer] = None,
) -> List[SlideTranslation]:
    """Collect translation data for review file generation."""
    review_slides = []
    
    try:
        # Parse translated XML
        tree = ET.parse(translated_xml_path)
        root = tree.getroot()
        
        # Parse original XML if available
        original_root = None
        if original_xml_path and original_xml_path.exists():
            original_tree = ET.parse(original_xml_path)
            original_root = original_tree.getroot()
        
        for xml_slide in root.findall(".//slide"):
            slide_number = int(xml_slide.get("number", "0"))
            if slide_number == 0:
                continue
            
            original_texts = []
            translated_texts = []
            quality_score = None
            issues = []
            
            # Collect text elements
            for text_elem in xml_slide.findall(".//text_element"):
                shape_index = int(text_elem.get("shape_index", "-1"))
                props_elem = text_elem.find("properties")
                if props_elem is not None and props_elem.text:
                    try:
                        shape_data = json.loads(props_elem.text)
                        translated_texts.append({
                            "shape_index": shape_index,
                            "text": shape_data.get("text", ""),
                            "properties": shape_data,
                        })
                    except Exception:
                        pass
            
            # Get original texts if available
            if original_root:
                orig_slide = original_root.find(f".//slide[@number='{slide_number}']")
                if orig_slide is not None:
                    for text_elem in orig_slide.findall(".//text_element"):
                        shape_index = int(text_elem.get("shape_index", "-1"))
                        props_elem = text_elem.find("properties")
                        if props_elem is not None and props_elem.text:
                            try:
                                shape_data = json.loads(props_elem.text)
                                original_texts.append({
                                    "shape_index": shape_index,
                                    "text": shape_data.get("text", ""),
                                    "properties": shape_data,
                                })
                            except Exception:
                                pass
            
            # Match original and translated texts by shape_index
            # Create a map for easier matching
            orig_map = {t["shape_index"]: t for t in original_texts}
            trans_map = {t["shape_index"]: t for t in translated_texts}
            
            # Ensure both lists are aligned
            all_indices = sorted(set(orig_map.keys()) | set(trans_map.keys()))
            aligned_original = []
            aligned_translated = []
            
            for idx in all_indices:
                aligned_original.append(orig_map.get(idx, {"shape_index": idx, "text": "", "properties": {}}))
                aligned_translated.append(trans_map.get(idx, {"shape_index": idx, "text": "", "properties": {}}))
            
            slide_translation = SlideTranslation(
                slide_number=slide_number,
                original_texts=aligned_original,
                translated_texts=aligned_translated,
                quality_score=quality_score,
                issues=issues,
            )
            review_slides.append(slide_translation)
    
    except Exception as e:
        print(f"Warning: Could not collect review data: {e}")
    
    return review_slides


def regenerate_ppt_from_review(
    original_ppt_path: Path,
    review_data,
    file_index: int,
    output_path: Path,
    update_memory=None,
) -> None:
    """Regenerate PPT from edited review file data.
    
    Args:
        original_ppt_path: Path to original PPT file
        review_data: ReviewData object with edited translations
        file_index: Index of file in review_data.files list
        output_path: Path for output PPT file
        update_memory: Optional TranslationMemory to update with edits
    """
    if file_index >= len(review_data.files):
        raise ValueError(f"File index {file_index} out of range")
    
    file_review = review_data.files[file_index]
    
    try:
        prs = Presentation(str(original_ppt_path))
        
        # Create a mapping of slide number to translations
        slide_translations = {s.slide_number: s for s in file_review.slides}
        
        for slide_number, slide in enumerate(prs.slides, start=1):
            if slide_number not in slide_translations:
                continue
            
            slide_trans = slide_translations[slide_number]
            
            # Create a mapping of shape_index to translated text
            trans_map = {t["shape_index"]: t for t in slide_trans.translated_texts}
            
            for shape_index, shape in enumerate(slide.shapes):
                if shape_index in trans_map:
                    trans_data = trans_map[shape_index]
                    
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        # Handle table updates if needed
                        pass
                    elif hasattr(shape, "text"):
                        # Update text with edited translation
                        edited_text = trans_data.get("text", "")
                        if edited_text:
                            # Get properties from original or use defaults
                            props = trans_data.get("properties", {})
                            
                            # Clear existing text
                            shape.text = ""
                            paragraph = shape.text_frame.paragraphs[0]
                            run = paragraph.add_run()
                            run.text = edited_text
                            
                            # Apply properties if available
                            if props.get("font_size"):
                                run.font.size = Pt(props["font_size"] * 0.7)
                            if props.get("font_name"):
                                run.font.name = props["font_name"]
                            if props.get("font_color"):
                                try:
                                    run.font.color.rgb = RGBColor.from_string(props["font_color"])
                                except Exception:
                                    pass
                            
                            # Update translation memory if provided
                            if update_memory:
                                orig_text = ""
                                for orig in slide_trans.original_texts:
                                    if orig.get("shape_index") == shape_index:
                                        orig_text = orig.get("text", "")
                                        break
                                if orig_text:
                                    update_memory.set(orig_text, edited_text)
        
        prs.save(str(output_path))
        print(f"✅ Regenerated PPT from review: {output_path.name}")
    
    except Exception as e:
        raise RuntimeError(f"Error regenerating PPT: {e}") from e


def cleanup_intermediate_files(temp_dir: Path) -> None:
    """Remove all intermediate files and the temp directory."""
    try:
        # Remove all files in the temp directory
        for file in temp_dir.glob("*"):
            if file.is_file():
                try:
                    file.unlink()
                except Exception:
                    pass  # Continue cleaning up other files even if one fails
        # Remove the temp directory itself (only works if empty)
        if temp_dir.exists():
            try:
                temp_dir.rmdir()
            except OSError:
                # Directory not empty or other error - log but don't fail
                print(f"Warning: Could not remove temp directory {temp_dir.name}/ (may not be empty)")
    except Exception as exc:  # pragma: no cover - logging only
        print(f"Warning: Could not clean up intermediate files: {exc}")


def process_ppt_file(
    ppt_path: Path,
    *,
    translator: TranslationService,
    source_lang: str,
    target_lang: str,
    max_workers: int = 4,
    cleanup: bool = True,
    vision_reviewer: VisionReviewer | None = None,
    max_refinement_iterations: int = 3,
    collect_review_data: bool = False,
):
    """Process a single PowerPoint file from extraction to translated output."""
    if not ppt_path.is_file():
        raise FileNotFoundError(f"'{ppt_path}' is not a valid file.")
    if ppt_path.suffix.lower() not in {".ppt", ".pptx"}:
        raise ValueError(f"'{ppt_path}' is not a PowerPoint file.")

    base_dir = ppt_path.parent
    # Create temp directory for intermediate files
    temp_dir = base_dir / f"{ppt_path.stem}_temp"
    temp_dir.mkdir(exist_ok=True)
    print(f"Using temp directory: {temp_dir.name}/")

    print(f"Generating original XML for {ppt_path.name}...")
    original_xml = ppt_to_xml(
        str(ppt_path),
        translator=None,
        source_lang=source_lang,
        target_lang=target_lang,
        max_workers=max_workers,
        temp_dir=temp_dir,
    )
    if original_xml:
        original_output_path = temp_dir / f"{ppt_path.stem}_original.xml"
        with open(original_output_path, "w", encoding="utf-8") as handle:
            handle.write(original_xml)
        print(f"Original XML saved: {original_output_path}")

    print(
        f"Generating translated XML (from {source_lang} to {target_lang}) for {ppt_path.name}..."
    )
    translated_xml = ppt_to_xml(
        str(ppt_path),
        translator=translator,
        source_lang=source_lang,
        target_lang=target_lang,
        max_workers=max_workers,
        temp_dir=temp_dir,
    )
    if not translated_xml:
        return None, None if collect_review_data else None

    translated_output_path = temp_dir / f"{ppt_path.stem}_translated.xml"
    with open(translated_output_path, "w", encoding="utf-8") as handle:
        handle.write(translated_xml)
    print(f"Translated XML saved: {translated_output_path}")

    print(f"Creating translated PPT for {ppt_path.name}...")
    output_filename = f"{ppt_path.stem}_translated{ppt_path.suffix}"
    output_ppt_path = base_dir / output_filename
    
    # Vision-based iterative refinement
    if vision_reviewer:
        print(f"🔍 Vision review enabled (max {max_refinement_iterations} iterations)")
        
        # Pre-translation: Analyze original slides
        print("  📸 Rendering original slides for analysis...")
        prs_original = Presentation(str(ppt_path))
        for slide_num in range(1, len(prs_original.slides) + 1):
            original_img = temp_dir / f"slide_{slide_num}_original.png"
            if render_slide_to_image(ppt_path, slide_num, original_img):
                analysis = vision_reviewer.analyze_original_slide(
                    original_img, source_lang, target_lang
                )
                if analysis:
                    print(f"    Slide {slide_num}: Analyzed")
        
        # Iterative refinement loop
        best_ppt_path = None
        best_score = 0.0
        
        for iteration in range(1, max_refinement_iterations + 1):
            if iteration > 1:
                print(f"  🔄 Refinement iteration {iteration}/{max_refinement_iterations}")
            
            # Create translated PPTX
            iter_output_path = base_dir / f"{ppt_path.stem}_translated_iter{iteration}{ppt_path.suffix}"
            create_translated_ppt(str(ppt_path), str(translated_output_path), str(iter_output_path))
            
            # Review translated slides
            print(f"  📸 Rendering translated slides for review...")
            all_passed = True
            min_score = 10.0
            
            for slide_num in range(1, len(prs_original.slides) + 1):
                original_img = temp_dir / f"slide_{slide_num}_original.png"
                translated_img = temp_dir / f"slide_{slide_num}_translated_iter{iteration}.png"
                
                if render_slide_to_image(iter_output_path, slide_num, translated_img):
                    review_result = vision_reviewer.review_translated_slide(
                        original_img, translated_img, source_lang, target_lang
                    )
                    
                    print(f"    Slide {slide_num}: Quality score {review_result.quality_score:.1f}/10")
                    if review_result.issues:
                        print(f"      Issues: {', '.join(review_result.issues[:3])}")
                    
                    min_score = min(min_score, review_result.quality_score)
                    if review_result.needs_refinement:
                        all_passed = False
            
            # Check if quality threshold met
            if min_score >= vision_reviewer.quality_threshold or all_passed:
                print(f"  ✅ Quality threshold met (score: {min_score:.1f}/10)")
                best_ppt_path = iter_output_path
                break
            elif iteration < max_refinement_iterations:
                print(f"  ⚠️  Quality below threshold ({min_score:.1f}/10), refining...")
                # TODO: Apply suggestions from review_result to improve translation
                # For now, we'll just retry (could use suggestions to adjust font sizes, etc.)
            else:
                print(f"  ⚠️  Max iterations reached (final score: {min_score:.1f}/10)")
                best_ppt_path = iter_output_path
        
        # Use best result or final iteration
        if best_ppt_path and best_ppt_path != output_ppt_path:
            best_ppt_path.rename(output_ppt_path)
            # Clean up iteration files
            for iter_file in base_dir.glob(f"{ppt_path.stem}_translated_iter*.pptx"):
                if iter_file != output_ppt_path:
                    iter_file.unlink()
    else:
        # No vision review - just create translated PPT
        create_translated_ppt(str(ppt_path), str(translated_output_path), str(output_ppt_path))

    # Collect review data if requested
    review_slides = None
    if collect_review_data:
        review_slides = _collect_review_data(
            original_xml_path=original_output_path if original_xml else None,
            translated_xml_path=translated_output_path,
            vision_reviewer=vision_reviewer,
        )

    if cleanup:
        cleanup_intermediate_files(temp_dir)
        print(f"Intermediate files and temp directory cleaned up.")
    else:
        print(f"Intermediate files kept in {temp_dir.name}/ (you can delete manually)")

    if collect_review_data:
        return output_ppt_path, review_slides
    return output_ppt_path
