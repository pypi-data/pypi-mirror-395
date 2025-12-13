"""Slide rendering utilities for vision-based review."""
from __future__ import annotations

import subprocess
import tempfile
from pathlib import Path
from typing import Optional

from pptx import Presentation


def render_slide_to_image(
    ppt_path: Path, slide_number: int, output_path: Path, width: int = 1920, height: int = 1080
) -> bool:
    """Render a specific slide from PPT to PNG image.
    
    Tries multiple methods:
    1. LibreOffice headless (best quality, cross-platform)
    2. Pillow-based simple rendering (fallback)
    
    Returns True if successful.
    """
    try:
        # Method 1: Try LibreOffice (best quality)
        if _render_with_libreoffice(ppt_path, slide_number, output_path):
            return True
        
        # Method 2: Fallback to Pillow-based rendering
        if _render_with_pillow(ppt_path, slide_number, output_path, width, height):
            return True
        
        return False
    except Exception as e:
        print(f"Warning: Could not render slide {slide_number}: {e}")
        return False


def _render_with_libreoffice(ppt_path: Path, slide_number: int, output_path: Path) -> bool:
    """Render using LibreOffice headless mode."""
    try:
        # Create single-slide presentation
        prs = Presentation(str(ppt_path))
        if slide_number > len(prs.slides):
            return False
        
        slide = prs.slides[slide_number - 1]
        
        # Create temp single-slide PPTX
        temp_prs = Presentation()
        temp_prs.slide_width = prs.slide_width
        temp_prs.slide_height = prs.slide_height
        
        # Copy slide layout
        layout = temp_prs.slide_layouts[0]  # Use blank layout
        temp_slide = temp_prs.slides.add_slide(layout)
        
        # Copy shapes (simplified - just copy text for now)
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                textbox = temp_slide.shapes.add_textbox(
                    shape.left, shape.top, shape.width, shape.height
                )
                textbox.text = shape.text
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            temp_pptx = Path(tmp.name)
            temp_prs.save(str(temp_pptx))
        
        try:
            # Try LibreOffice conversion
            result = subprocess.run(
                [
                    'libreoffice',
                    '--headless',
                    '--convert-to', 'png',
                    '--outdir', str(output_path.parent),
                    str(temp_pptx)
                ],
                capture_output=True,
                timeout=30,
                check=False,
            )
            
            if result.returncode == 0:
                png_path = output_path.parent / f"{temp_pptx.stem}.png"
                if png_path.exists():
                    png_path.rename(output_path)
                    temp_pptx.unlink()
                    return True
        except (FileNotFoundError, subprocess.TimeoutExpired):
            pass
        
        temp_pptx.unlink()
        return False
    except Exception:
        return False


def _render_with_pillow(ppt_path: Path, slide_number: int, output_path: Path, width: int, height: int) -> bool:
    """Fallback rendering using Pillow (simplified visual representation)."""
    try:
        from PIL import Image, ImageDraw, ImageFont
        
        prs = Presentation(str(ppt_path))
        if slide_number > len(prs.slides):
            return False
        
        slide = prs.slides[slide_number - 1]
        
        # Create image
        img = Image.new('RGB', (width, height), color='white')
        draw = ImageDraw.Draw(img)
        
        # Try to load a font
        try:
            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
        except:
            font = ImageFont.load_default()
        
        # Draw text elements
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                # Convert EMU to pixels
                x = int(shape.left / 914400 * width)
                y = int(shape.top / 914400 * height)
                w = int(shape.width / 914400 * width)
                h = int(shape.height / 914400 * height)
                
                # Draw text (simplified - doesn't handle all formatting)
                draw.text((x, y), shape.text, fill='black', font=font)
        
        img.save(output_path, 'PNG')
        return True
    except ImportError:
        print("Warning: Pillow not installed. Install with: pip install Pillow")
        return False
    except Exception:
        return False

